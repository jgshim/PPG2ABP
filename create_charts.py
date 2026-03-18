"""
Generate comprehensive result charts and update PPT with them.
"""
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import os

plt.rcParams['font.family'] = 'Malgun Gothic'
plt.rcParams['axes.unicode_minus'] = False
plt.rcParams['figure.dpi'] = 200

BASE = 'C:/Users/jaege/Desktop/Study/PPG2ABP'
FIG_DIR = os.path.join(BASE, 'figures')
os.makedirs(FIG_DIR, exist_ok=True)

# ======================================================================
# DATA
# ======================================================================
models = ['XGBoost', 'LightGBM', 'Light CNN\n(CPU)', 'Improved CNN\n(GPU)', 'ResNet1D\n(GPU)']
models_short = ['XGBoost', 'LightGBM', 'Light CNN', 'Imp. CNN', 'ResNet1D']

sbp_mae = [14.98, 16.09, 12.95, 12.64, 12.25]
dbp_mae = [8.66, 8.42, 7.47, 7.20, 7.83]
mbp_mae = [10.41, 10.15, 8.49, 8.08, 8.22]

sbp_rmse = [19.47, 20.62, 16.76, 16.21, 15.91]
dbp_rmse = [11.62, 11.33, 10.21, 9.83, 10.18]
mbp_rmse = [13.96, 13.65, 11.67, 11.14, 11.11]

sbp_r2 = [-0.1925, -0.3375, 0.1569, 0.1731, 0.2031]
dbp_r2 = [-0.0780, -0.0249, 0.2046, 0.2282, 0.1724]
mbp_r2 = [-0.1436, -0.0935, 0.2193, 0.2717, 0.2763]

# BHS cumulative error % (only for CNN/ResNet GPU)
bhs_models = ['Imp. CNN', 'ResNet1D']
cnn_sbp_bhs = [25.3, 48.1, 66.9]
cnn_dbp_bhs = [44.8, 76.7, 90.0]
cnn_mbp_bhs = [42.5, 70.8, 86.1]
res_sbp_bhs = [27.5, 51.9, 67.9]
res_dbp_bhs = [39.0, 71.7, 89.2]
res_mbp_bhs = [40.2, 70.5, 85.0]

# Bias +/- SD
bias_models = ['Imp. CNN', 'ResNet1D']
cnn_bias = [-0.9, -0.7, -0.9]
cnn_sd = [16.2, 9.8, 11.1]
res_bias = [-2.2, -0.5, -0.8]
res_sd = [15.8, 10.2, 11.1]

# Fold-by-fold
folds = [1, 2, 3, 4, 5]
cnn_fold_sbp = [12.10, 12.36, 11.31, 13.66, 13.77]
cnn_fold_dbp = [7.72, 6.65, 6.66, 7.89, 7.08]
cnn_fold_mbp = [8.89, 7.65, 7.36, 8.57, 7.93]
res_fold_sbp = [14.02, 10.57, 11.75, 12.01, 12.91]
res_fold_dbp = [9.09, 7.49, 7.51, 7.28, 7.76]
res_fold_mbp = [9.31, 7.78, 8.31, 7.84, 7.87]

cnn_fold_epochs = [18, 17, 25, 17, 10]
res_fold_epochs = [10, 22, 13, 23, 13]
cnn_fold_time = [411, 388, 570, 380, 231]
res_fold_time = [316, 677, 410, 701, 409]

# Colors
C_ML = ['#7f8c8d', '#95a5a6']  # grays for ML
C_DL = ['#3498db', '#2980b9', '#e74c3c']  # blues for CNN, red for ResNet
COLORS5 = ['#7f8c8d', '#95a5a6', '#3498db', '#2980b9', '#e74c3c']
C_SBP = '#e74c3c'
C_DBP = '#3498db'
C_MBP = '#2ecc71'


# ======================================================================
# CHART 1: MAE Comparison Bar Chart (Grouped)
# ======================================================================
fig, ax = plt.subplots(figsize=(10, 5.5))
x = np.arange(len(models_short))
w = 0.25
bars1 = ax.bar(x - w, sbp_mae, w, label='SBP', color=C_SBP, alpha=0.85, edgecolor='white', linewidth=0.5)
bars2 = ax.bar(x, dbp_mae, w, label='DBP', color=C_DBP, alpha=0.85, edgecolor='white', linewidth=0.5)
bars3 = ax.bar(x + w, mbp_mae, w, label='MBP', color=C_MBP, alpha=0.85, edgecolor='white', linewidth=0.5)

for bars in [bars1, bars2, bars3]:
    for bar in bars:
        h = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2, h + 0.15, f'{h:.1f}',
                ha='center', va='bottom', fontsize=8, fontweight='bold')

ax.set_ylabel('MAE (mmHg)', fontsize=12, fontweight='bold')
ax.set_title('Model Comparison - MAE (SBP / DBP / MBP)', fontsize=14, fontweight='bold', pad=10)
ax.set_xticks(x)
ax.set_xticklabels(models_short, fontsize=11)
ax.legend(fontsize=11, loc='upper right')
ax.set_ylim(0, 20)
ax.axhline(y=5, color='green', linestyle='--', alpha=0.4, linewidth=1)
ax.text(4.5, 5.2, 'AAMI Target (5 mmHg)', fontsize=8, color='green', alpha=0.6)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.grid(axis='y', alpha=0.3)
# Highlight best
ax.annotate('Best SBP', xy=(4, 12.25), xytext=(4.3, 14.5),
            arrowprops=dict(arrowstyle='->', color=C_SBP, lw=1.5),
            fontsize=9, color=C_SBP, fontweight='bold')
ax.annotate('Best DBP', xy=(3, 7.20), xytext=(3.3, 4.5),
            arrowprops=dict(arrowstyle='->', color=C_DBP, lw=1.5),
            fontsize=9, color=C_DBP, fontweight='bold')
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart1_mae_comparison.png'), bbox_inches='tight')
plt.close()
print("Chart 1: MAE comparison saved")


# ======================================================================
# CHART 2: R-squared Comparison
# ======================================================================
fig, ax = plt.subplots(figsize=(10, 5.5))
x = np.arange(len(models_short))
w = 0.25
bars1 = ax.bar(x - w, sbp_r2, w, label='SBP', color=C_SBP, alpha=0.85, edgecolor='white')
bars2 = ax.bar(x, dbp_r2, w, label='DBP', color=C_DBP, alpha=0.85, edgecolor='white')
bars3 = ax.bar(x + w, mbp_r2, w, label='MBP', color=C_MBP, alpha=0.85, edgecolor='white')

for bars in [bars1, bars2, bars3]:
    for bar in bars:
        h = bar.get_height()
        ypos = h + 0.01 if h >= 0 else h - 0.03
        va = 'bottom' if h >= 0 else 'top'
        ax.text(bar.get_x() + bar.get_width()/2, ypos, f'{h:.2f}',
                ha='center', va=va, fontsize=7.5, fontweight='bold')

ax.set_ylabel(r'$R^2$', fontsize=13, fontweight='bold')
ax.set_title(r'Model Comparison - $R^2$ Score (higher is better)', fontsize=14, fontweight='bold', pad=10)
ax.set_xticks(x)
ax.set_xticklabels(models_short, fontsize=11)
ax.legend(fontsize=11, loc='lower right')
ax.axhline(y=0, color='black', linestyle='-', alpha=0.3, linewidth=0.8)
ax.set_ylim(-0.45, 0.35)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.grid(axis='y', alpha=0.3)
# Add zone labels
ax.axhspan(-0.45, 0, alpha=0.05, color='red')
ax.axhspan(0, 0.35, alpha=0.05, color='green')
ax.text(4.6, -0.05, 'Worse than mean', fontsize=8, color='red', alpha=0.5, ha='right')
ax.text(4.6, 0.01, 'Better than mean', fontsize=8, color='green', alpha=0.5, ha='right')
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart2_r2_comparison.png'), bbox_inches='tight')
plt.close()
print("Chart 2: R2 comparison saved")


# ======================================================================
# CHART 3: BHS Cumulative Error Distribution
# ======================================================================
fig, axes = plt.subplots(1, 3, figsize=(13, 5))
thresholds = ['<=5 mmHg', '<=10 mmHg', '<=15 mmHg']

# BHS grade thresholds
bhs_a = [60, 85, 95]
bhs_b = [50, 75, 90]
bhs_c = [40, 65, 85]

for idx, (ax, target, cnn_vals, res_vals) in enumerate(zip(
    axes,
    ['SBP', 'DBP', 'MBP'],
    [cnn_sbp_bhs, cnn_dbp_bhs, cnn_mbp_bhs],
    [res_sbp_bhs, res_dbp_bhs, res_mbp_bhs]
)):
    x = np.arange(3)
    w = 0.3
    bars1 = ax.bar(x - w/2, cnn_vals, w, label='Imp. CNN', color='#2980b9', alpha=0.85, edgecolor='white')
    bars2 = ax.bar(x + w/2, res_vals, w, label='ResNet1D', color='#e74c3c', alpha=0.85, edgecolor='white')

    for bars in [bars1, bars2]:
        for bar in bars:
            h = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2, h + 0.8, f'{h:.1f}%',
                    ha='center', va='bottom', fontsize=8, fontweight='bold')

    # BHS grade lines
    for grade, vals, color, ls in [
        ('A', bhs_a, '#27ae60', '-'),
        ('B', bhs_b, '#f39c12', '--'),
        ('C', bhs_c, '#e67e22', ':'),
    ]:
        ax.plot(x, [vals[i] for i in range(3)], color=color, linestyle=ls,
                marker='_', markersize=10, linewidth=1.5, alpha=0.6)
        ax.text(2.35, vals[2], f'Grade {grade}', fontsize=7, color=color, va='center', fontweight='bold')

    ax.set_title(f'{target}', fontsize=14, fontweight='bold',
                 color={'SBP': C_SBP, 'DBP': C_DBP, 'MBP': C_MBP}[target])
    ax.set_xticks(x)
    ax.set_xticklabels(thresholds, fontsize=10)
    ax.set_ylim(0, 105)
    ax.set_ylabel('Cumulative %' if idx == 0 else '', fontsize=11)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.3)
    if idx == 0:
        ax.legend(fontsize=9, loc='upper left')

fig.suptitle('BHS Cumulative Error Distribution', fontsize=15, fontweight='bold', y=1.02)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart3_bhs_distribution.png'), bbox_inches='tight')
plt.close()
print("Chart 3: BHS distribution saved")


# ======================================================================
# CHART 4: Bias +/- SD (Bland-Altman Summary)
# ======================================================================
fig, ax = plt.subplots(figsize=(9, 5.5))
targets = ['SBP', 'DBP', 'MBP']
y_pos = np.array([0, 1, 2])
offset = 0.15

# CNN
for i, (bias, sd, target) in enumerate(zip(cnn_bias, cnn_sd, targets)):
    ax.errorbar(bias, y_pos[i] - offset, xerr=sd, fmt='o', markersize=8,
                color='#2980b9', capsize=6, capthick=2, linewidth=2, label='Imp. CNN' if i == 0 else None)
    ax.text(bias + sd + 0.5, y_pos[i] - offset, f'{bias:+.1f} +/- {sd:.1f}',
            fontsize=9, va='center', color='#2980b9', fontweight='bold')

# ResNet
for i, (bias, sd, target) in enumerate(zip(res_bias, res_sd, targets)):
    ax.errorbar(bias, y_pos[i] + offset, xerr=sd, fmt='s', markersize=8,
                color='#e74c3c', capsize=6, capthick=2, linewidth=2, label='ResNet1D' if i == 0 else None)
    ax.text(bias + sd + 0.5, y_pos[i] + offset, f'{bias:+.1f} +/- {sd:.1f}',
            fontsize=9, va='center', color='#e74c3c', fontweight='bold')

ax.axvline(x=0, color='black', linestyle='-', linewidth=0.8, alpha=0.5)
# AAMI limits
ax.axvline(x=5, color='green', linestyle='--', linewidth=1, alpha=0.3)
ax.axvline(x=-5, color='green', linestyle='--', linewidth=1, alpha=0.3)
ax.axvspan(-8, 8, alpha=0.04, color='green')
ax.text(6, 2.5, 'AAMI: |bias|<=5, SD<=8', fontsize=8, color='green', alpha=0.6)

ax.set_yticks(y_pos)
ax.set_yticklabels(targets, fontsize=13, fontweight='bold')
ax.set_xlabel('Bias +/- SD (mmHg)', fontsize=12, fontweight='bold')
ax.set_title('Bland-Altman Summary: Bias +/- SD', fontsize=14, fontweight='bold', pad=10)
ax.legend(fontsize=11, loc='lower left')
ax.set_xlim(-25, 25)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.grid(axis='x', alpha=0.3)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart4_bias_sd.png'), bbox_inches='tight')
plt.close()
print("Chart 4: Bias+/-SD saved")


# ======================================================================
# CHART 5: Fold-by-Fold MAE Comparison (CNN vs ResNet)
# ======================================================================
fig, axes = plt.subplots(1, 3, figsize=(14, 5))

for ax, target, cnn_vals, res_vals in zip(
    axes,
    ['SBP MAE', 'DBP MAE', 'MBP MAE'],
    [cnn_fold_sbp, cnn_fold_dbp, cnn_fold_mbp],
    [res_fold_sbp, res_fold_dbp, res_fold_mbp]
):
    ax.plot(folds, cnn_vals, 'o-', color='#2980b9', linewidth=2, markersize=8, label='Imp. CNN')
    ax.plot(folds, res_vals, 's--', color='#e74c3c', linewidth=2, markersize=8, label='ResNet1D')

    # Mean lines
    cnn_mean = np.mean(cnn_vals)
    res_mean = np.mean(res_vals)
    ax.axhline(y=cnn_mean, color='#2980b9', linestyle=':', alpha=0.5, linewidth=1)
    ax.axhline(y=res_mean, color='#e74c3c', linestyle=':', alpha=0.5, linewidth=1)
    ax.text(5.15, cnn_mean, f'{cnn_mean:.2f}', fontsize=8, color='#2980b9', va='center')
    ax.text(5.15, res_mean, f'{res_mean:.2f}', fontsize=8, color='#e74c3c', va='center')

    # Values on points
    for i in range(5):
        ax.text(folds[i], cnn_vals[i] - 0.3, f'{cnn_vals[i]:.1f}', ha='center', fontsize=7, color='#2980b9')
        ax.text(folds[i], res_vals[i] + 0.25, f'{res_vals[i]:.1f}', ha='center', fontsize=7, color='#e74c3c')

    color = {'SBP MAE': C_SBP, 'DBP MAE': C_DBP, 'MBP MAE': C_MBP}[target]
    ax.set_title(target, fontsize=14, fontweight='bold', color=color)
    ax.set_xlabel('Fold', fontsize=11)
    ax.set_ylabel('MAE (mmHg)', fontsize=11)
    ax.set_xticks(folds)
    ax.legend(fontsize=9)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(alpha=0.3)

fig.suptitle('Fold-by-Fold Performance Comparison', fontsize=15, fontweight='bold', y=1.02)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart5_fold_comparison.png'), bbox_inches='tight')
plt.close()
print("Chart 5: Fold comparison saved")


# ======================================================================
# CHART 6: Radar/Spider Chart - Multi-Metric Comparison
# ======================================================================
fig, ax = plt.subplots(figsize=(7, 7), subplot_kw=dict(polar=True))

# Metrics: SBP_MAE(inv), DBP_MAE(inv), MBP_MAE(inv), SBP_R2, DBP_R2, MBP_R2
# Normalize to 0-1 where 1 is best
def normalize(val, worst, best):
    return max(0, min(1, (val - worst) / (best - worst)))

categories = ['SBP MAE\n(lower=better)', 'DBP MAE\n(lower=better)', 'MBP MAE\n(lower=better)',
              r'SBP $R^2$', r'DBP $R^2$', r'MBP $R^2$']
N = len(categories)

# Range for normalization
sbp_mae_range = (20, 10)  # worst, best (inverted)
dbp_mae_range = (12, 5)
mbp_mae_range = (13, 6)
r2_range = (-0.4, 0.5)

model_radar = {
    'XGBoost': [normalize(14.98, *sbp_mae_range), normalize(8.66, *dbp_mae_range),
                normalize(10.41, *mbp_mae_range), normalize(-0.19, *r2_range),
                normalize(-0.08, *r2_range), normalize(-0.14, *r2_range)],
    'LightGBM': [normalize(16.09, *sbp_mae_range), normalize(8.42, *dbp_mae_range),
                 normalize(10.15, *mbp_mae_range), normalize(-0.34, *r2_range),
                 normalize(-0.02, *r2_range), normalize(-0.09, *r2_range)],
    'Light CNN': [normalize(12.95, *sbp_mae_range), normalize(7.47, *dbp_mae_range),
                  normalize(8.49, *mbp_mae_range), normalize(0.16, *r2_range),
                  normalize(0.20, *r2_range), normalize(0.22, *r2_range)],
    'Imp. CNN': [normalize(12.64, *sbp_mae_range), normalize(7.20, *dbp_mae_range),
                 normalize(8.08, *mbp_mae_range), normalize(0.17, *r2_range),
                 normalize(0.23, *r2_range), normalize(0.27, *r2_range)],
    'ResNet1D': [normalize(12.25, *sbp_mae_range), normalize(7.83, *dbp_mae_range),
                 normalize(8.22, *mbp_mae_range), normalize(0.20, *r2_range),
                 normalize(0.17, *r2_range), normalize(0.28, *r2_range)],
}

angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
angles += angles[:1]

colors_radar = {'XGBoost': '#7f8c8d', 'LightGBM': '#95a5a6',
                'Light CNN': '#3498db', 'Imp. CNN': '#2980b9', 'ResNet1D': '#e74c3c'}
lw_radar = {'XGBoost': 1.2, 'LightGBM': 1.2, 'Light CNN': 1.8, 'Imp. CNN': 2.5, 'ResNet1D': 2.5}

for name, vals in model_radar.items():
    values = vals + vals[:1]
    ax.plot(angles, values, 'o-', linewidth=lw_radar[name], label=name,
            color=colors_radar[name], markersize=4 if 'GB' in name else 6)
    if name in ['Imp. CNN', 'ResNet1D']:
        ax.fill(angles, values, alpha=0.08, color=colors_radar[name])

ax.set_xticks(angles[:-1])
ax.set_xticklabels(categories, fontsize=9)
ax.set_ylim(0, 1)
ax.set_yticks([0.2, 0.4, 0.6, 0.8])
ax.set_yticklabels(['0.2', '0.4', '0.6', '0.8'], fontsize=8, alpha=0.5)
ax.set_title('Multi-Metric Radar Comparison\n(outer = better)', fontsize=13, fontweight='bold', pad=20)
ax.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1), fontsize=9)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart6_radar.png'), bbox_inches='tight')
plt.close()
print("Chart 6: Radar chart saved")


# ======================================================================
# CHART 7: Performance Heatmap
# ======================================================================
fig, ax = plt.subplots(figsize=(10, 5))

metrics = ['SBP MAE', 'DBP MAE', 'MBP MAE', 'SBP RMSE', 'DBP RMSE', 'MBP RMSE']
data = np.array([
    [14.98, 8.66, 10.41, 19.47, 11.62, 13.96],  # XGBoost
    [16.09, 8.42, 10.15, 20.62, 11.33, 13.65],  # LightGBM
    [12.95, 7.47, 8.49, 16.76, 10.21, 11.67],   # Light CNN
    [12.64, 7.20, 8.08, 16.21, 9.83, 11.14],    # Imp. CNN
    [12.25, 7.83, 8.22, 15.91, 10.18, 11.11],   # ResNet1D
])

im = ax.imshow(data, cmap='RdYlGn_r', aspect='auto')
ax.set_xticks(np.arange(len(metrics)))
ax.set_yticks(np.arange(len(models_short)))
ax.set_xticklabels(metrics, fontsize=11, fontweight='bold')
ax.set_yticklabels(models_short, fontsize=11, fontweight='bold')

for i in range(len(models_short)):
    for j in range(len(metrics)):
        val = data[i, j]
        # Find min in column
        col_min = data[:, j].min()
        weight = 'bold' if val == col_min else 'normal'
        color = 'white' if val > 14 else 'black'
        ax.text(j, i, f'{val:.2f}', ha='center', va='center',
                fontsize=10, fontweight=weight, color=color)

ax.set_title('Performance Heatmap (lower = better, bold = best)', fontsize=14, fontweight='bold', pad=10)
cbar = plt.colorbar(im, ax=ax, shrink=0.8)
cbar.set_label('Error (mmHg)', fontsize=11)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart7_heatmap.png'), bbox_inches='tight')
plt.close()
print("Chart 7: Heatmap saved")


# ======================================================================
# CHART 8: Training Efficiency (Time vs Performance)
# ======================================================================
fig, ax = plt.subplots(figsize=(9, 6))

# Models with time info
eff_models = ['Light CNN\n(CPU, sub)', 'Imp. CNN\n(GPU, full)', 'ResNet1D\n(GPU, full)']
eff_time_min = [5, 33, 42]  # approximate minutes
eff_dbp_mae = [7.47, 7.20, 7.83]
eff_sbp_mae = [12.95, 12.64, 12.25]
eff_mbp_mae = [8.49, 8.08, 8.22]
eff_data_size = [22638, 81812, 81812]
colors_eff = ['#3498db', '#2980b9', '#e74c3c']

# Bubble size based on data size
sizes = [s/200 for s in eff_data_size]

# SBP
scatter1 = ax.scatter(eff_time_min, eff_sbp_mae, s=sizes, c=colors_eff, alpha=0.6,
                       edgecolors='black', linewidth=1, zorder=5)
# DBP
scatter2 = ax.scatter(eff_time_min, eff_dbp_mae, s=sizes, c=colors_eff, alpha=0.6,
                       edgecolors='black', linewidth=1, marker='s', zorder=5)

# Connect SBP-DBP pairs
for i in range(3):
    ax.plot([eff_time_min[i], eff_time_min[i]], [eff_sbp_mae[i], eff_dbp_mae[i]],
            color=colors_eff[i], linestyle=':', alpha=0.5, linewidth=1.5)

# Labels
labels = ['Light CNN', 'Imp. CNN', 'ResNet1D']
for i in range(3):
    ax.annotate(f'{labels[i]}\nSBP={eff_sbp_mae[i]:.2f}',
                (eff_time_min[i], eff_sbp_mae[i]),
                textcoords="offset points", xytext=(15, 10),
                fontsize=9, color=colors_eff[i], fontweight='bold')
    ax.annotate(f'DBP={eff_dbp_mae[i]:.2f}',
                (eff_time_min[i], eff_dbp_mae[i]),
                textcoords="offset points", xytext=(15, -15),
                fontsize=9, color=colors_eff[i], fontweight='bold')

# Legend
ax.scatter([], [], s=100, c='gray', alpha=0.5, label='SBP (circle)')
ax.scatter([], [], s=100, c='gray', alpha=0.5, marker='s', label='DBP (square)')
ax.text(0.98, 0.02, 'Bubble size = data size', transform=ax.transAxes,
        fontsize=8, ha='right', va='bottom', alpha=0.5, style='italic')

ax.set_xlabel('Training Time (minutes)', fontsize=12, fontweight='bold')
ax.set_ylabel('MAE (mmHg)', fontsize=12, fontweight='bold')
ax.set_title('Training Efficiency: Time vs. Performance', fontsize=14, fontweight='bold', pad=10)
ax.legend(fontsize=10, loc='upper right')
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.grid(alpha=0.3)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart8_efficiency.png'), bbox_inches='tight')
plt.close()
print("Chart 8: Efficiency saved")


# ======================================================================
# CHART 9: Model Progression (improvement over steps)
# ======================================================================
fig, axes = plt.subplots(1, 2, figsize=(12, 5))

steps = ['Step 2\nML\n(LOSO)', 'Step 3\nLight CNN\n(CPU)', 'Step 4\nImp. CNN\n(GPU)', 'Step 4\nResNet1D\n(GPU)']
# Best ML = XGBoost for SBP, LightGBM for DBP
prog_sbp = [14.98, 12.95, 12.64, 12.25]
prog_dbp = [8.42, 7.47, 7.20, 7.83]  # LGB best for ML
prog_mbp = [10.15, 8.49, 8.08, 8.22]

# Left: MAE progression
ax = axes[0]
x = np.arange(4)
ax.plot(x, prog_sbp, 'o-', color=C_SBP, linewidth=2.5, markersize=10, label='SBP')
ax.plot(x, prog_dbp, 's-', color=C_DBP, linewidth=2.5, markersize=10, label='DBP')
ax.plot(x, prog_mbp, '^-', color=C_MBP, linewidth=2.5, markersize=10, label='MBP')

for vals, color in [(prog_sbp, C_SBP), (prog_dbp, C_DBP), (prog_mbp, C_MBP)]:
    for i, v in enumerate(vals):
        ax.text(i, v + 0.25, f'{v:.2f}', ha='center', fontsize=9, color=color, fontweight='bold')

ax.fill_between(x, prog_sbp, prog_dbp, alpha=0.05, color='gray')
ax.set_xticks(x)
ax.set_xticklabels(steps, fontsize=9)
ax.set_ylabel('MAE (mmHg)', fontsize=12, fontweight='bold')
ax.set_title('MAE Improvement Across Steps', fontsize=13, fontweight='bold')
ax.legend(fontsize=10)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.grid(alpha=0.3)

# Right: Improvement %
ax = axes[1]
# % improvement from ML baseline
sbp_base, dbp_base, mbp_base = 14.98, 8.42, 10.15
imp_sbp = [(sbp_base - v) / sbp_base * 100 for v in prog_sbp]
imp_dbp = [(dbp_base - v) / dbp_base * 100 for v in prog_dbp]
imp_mbp = [(mbp_base - v) / mbp_base * 100 for v in prog_mbp]

w = 0.25
bars1 = ax.bar(x - w, imp_sbp, w, label='SBP', color=C_SBP, alpha=0.8)
bars2 = ax.bar(x, imp_dbp, w, label='DBP', color=C_DBP, alpha=0.8)
bars3 = ax.bar(x + w, imp_mbp, w, label='MBP', color=C_MBP, alpha=0.8)

for bars in [bars1, bars2, bars3]:
    for bar in bars:
        h = bar.get_height()
        if abs(h) > 1:
            ax.text(bar.get_x() + bar.get_width()/2, h + 0.3 if h >= 0 else h - 1.5,
                    f'{h:.1f}%', ha='center', fontsize=7.5, fontweight='bold')

ax.set_xticks(x)
ax.set_xticklabels(steps, fontsize=9)
ax.set_ylabel('Improvement from ML Baseline (%)', fontsize=11, fontweight='bold')
ax.set_title('Relative Improvement from ML Baseline', fontsize=13, fontweight='bold')
ax.legend(fontsize=10)
ax.axhline(y=0, color='black', linewidth=0.8, alpha=0.3)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.grid(axis='y', alpha=0.3)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart9_progression.png'), bbox_inches='tight')
plt.close()
print("Chart 9: Progression saved")


# ======================================================================
# CHART 10: U-Net Waveform Correlation
# ======================================================================
fig, ax = plt.subplots(figsize=(7, 5))

unet_folds = ['Fold 1', 'Fold 2', 'Fold 3', 'Mean\n(3 folds)']
unet_corr = [0.587, 0.772, 0.572, np.mean([0.587, 0.772, 0.572])]
unet_rmse = [13.40, 9.66, 14.03, np.mean([13.40, 9.66, 14.03])]
colors_unet = ['#27ae60', '#2ecc71', '#1abc9c', '#16a085']

ax2 = ax.twinx()

bars = ax.bar(np.arange(4), unet_corr, 0.5, color=colors_unet, alpha=0.8, edgecolor='white', linewidth=1.5)
line = ax2.plot(np.arange(4), unet_rmse, 'D-', color='#e74c3c', linewidth=2, markersize=10, label='WF RMSE')

for i, (c, r) in enumerate(zip(unet_corr, unet_rmse)):
    ax.text(i, c + 0.015, f'{c:.3f}', ha='center', fontsize=10, fontweight='bold', color='#1a5276')
    ax2.text(i, r + 0.3, f'{r:.2f}', ha='center', fontsize=9, fontweight='bold', color='#e74c3c')

ax.set_ylabel('Waveform Correlation', fontsize=12, fontweight='bold', color='#27ae60')
ax2.set_ylabel('Waveform RMSE (mmHg)', fontsize=12, fontweight='bold', color='#e74c3c')
ax.set_xticks(np.arange(4))
ax.set_xticklabels(unet_folds, fontsize=11)
ax.set_ylim(0, 1.0)
ax2.set_ylim(5, 18)
ax.set_title('U-Net 1D: Waveform Reconstruction Quality', fontsize=14, fontweight='bold', pad=10)

# Legend
bar_patch = mpatches.Patch(color='#27ae60', alpha=0.8, label='Correlation')
line_patch = plt.Line2D([0], [0], color='#e74c3c', marker='D', linewidth=2, label='WF RMSE')
ax.legend(handles=[bar_patch, line_patch], fontsize=10, loc='upper left')

ax.spines['top'].set_visible(False)
ax.grid(axis='y', alpha=0.2)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'chart10_unet_waveform.png'), bbox_inches='tight')
plt.close()
print("Chart 10: U-Net waveform saved")


# ======================================================================
# Now update PPTX with chart images
# ======================================================================
print("\n" + "="*60)
print("Inserting charts into PowerPoint...")
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation(os.path.join(BASE, 'PPG2ABP_Presentation.pptx'))

PRIMARY = PRGBColor(0, 51, 102)
WHITE = PRGBColor(255, 255, 255)

def add_chart_slide(title, img_path, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title bar
    shape = slide.shapes.add_shape(
        1, PInches(0), PInches(0), PInches(13.333), PInches(0.9))  # MSO_SHAPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(PInches(0.4), PInches(0.1), PInches(12), PInches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PPt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Malgun Gothic'

    # Image - centered, large
    if subtitle:
        slide.shapes.add_picture(img_path, PInches(0.5), PInches(1.1), PInches(12.3), PInches(5.8))
        txBox2 = slide.shapes.add_textbox(PInches(0.5), PInches(7.0), PInches(12), PInches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = PPt(12)
        p2.font.color.rgb = PRGBColor(100, 100, 100)
        p2.font.name = 'Malgun Gothic'
        p2.alignment = PP_ALIGN.CENTER
    else:
        slide.shapes.add_picture(img_path, PInches(0.5), PInches(1.1), PInches(12.3), PInches(6.2))
    return slide


def add_two_chart_slide(title, img1, img2, cap1='', cap2=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        1, PInches(0), PInches(0), PInches(13.333), PInches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(PInches(0.4), PInches(0.1), PInches(12), PInches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PPt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Malgun Gothic'

    slide.shapes.add_picture(img1, PInches(0.2), PInches(1.1), PInches(6.4), PInches(4.5))
    slide.shapes.add_picture(img2, PInches(6.7), PInches(1.1), PInches(6.4), PInches(4.5))

    if cap1:
        txBox1 = slide.shapes.add_textbox(PInches(0.2), PInches(5.7), PInches(6.4), PInches(0.3))
        p1 = txBox1.text_frame.paragraphs[0]
        p1.text = cap1; p1.font.size = PPt(11); p1.font.name = 'Malgun Gothic'
        p1.alignment = PP_ALIGN.CENTER; p1.font.color.rgb = PRGBColor(80, 80, 80)
    if cap2:
        txBox2 = slide.shapes.add_textbox(PInches(6.7), PInches(5.7), PInches(6.4), PInches(0.3))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = cap2; p2.font.size = PPt(11); p2.font.name = 'Malgun Gothic'
        p2.alignment = PP_ALIGN.CENTER; p2.font.color.rgb = PRGBColor(80, 80, 80)
    return slide


# Slide 1: MAE + R2 comparison
add_two_chart_slide(
    'Results - MAE & R2 Comparison',
    os.path.join(FIG_DIR, 'chart1_mae_comparison.png'),
    os.path.join(FIG_DIR, 'chart2_r2_comparison.png'),
    'MAE: Lower is better. Best SBP: ResNet1D (12.25), Best DBP: Imp. CNN (7.20)',
    'R2: Higher is better. DL models significantly outperform ML baselines'
)

# Slide 2: BHS Distribution
add_chart_slide(
    'Results - BHS Cumulative Error Distribution',
    os.path.join(FIG_DIR, 'chart3_bhs_distribution.png'),
    'BHS Grade A: >=60%/85%/95%, Grade B: >=50%/75%/90%, Grade C: >=40%/65%/85%'
)

# Slide 3: Bias+SD + Radar
add_two_chart_slide(
    'Results - Bias/SD & Multi-Metric Comparison',
    os.path.join(FIG_DIR, 'chart4_bias_sd.png'),
    os.path.join(FIG_DIR, 'chart6_radar.png'),
    'Bland-Altman Summary: Both models show low bias, but SD exceeds AAMI limit (8)',
    'Radar: DL models (blue/red) dominate ML baselines (gray) across all metrics'
)

# Slide 4: Fold-by-fold
add_chart_slide(
    'Results - Fold-by-Fold Performance',
    os.path.join(FIG_DIR, 'chart5_fold_comparison.png'),
    'CNN shows more stable fold performance; ResNet has higher variance especially in SBP'
)

# Slide 5: Heatmap + Efficiency
add_two_chart_slide(
    'Results - Performance Heatmap & Training Efficiency',
    os.path.join(FIG_DIR, 'chart7_heatmap.png'),
    os.path.join(FIG_DIR, 'chart8_efficiency.png'),
    'Green=better, Red=worse. Bold=best in column',
    'Bubble size = dataset size. GPU full-data models offer best accuracy'
)

# Slide 6: Progression + U-Net
add_two_chart_slide(
    'Results - Model Progression & U-Net Waveform',
    os.path.join(FIG_DIR, 'chart9_progression.png'),
    os.path.join(FIG_DIR, 'chart10_unet_waveform.png'),
    'Consistent improvement from ML to DL. Up to 18% SBP improvement',
    'U-Net waveform correlation: 0.644 avg (3 folds). Fold 2 best: 0.772'
)

# Save
pptx_path = os.path.join(BASE, 'PPG2ABP_Presentation.pptx')
try:
    prs.save(pptx_path)
    print(f"Saved: {pptx_path}")
except PermissionError:
    alt_path = os.path.join(BASE, 'PPG2ABP_Presentation_v2.pptx')
    prs.save(alt_path)
    print(f"Original locked. Saved as: {alt_path}")

print("\nAll done! 10 charts generated, 6 new slides added.")
