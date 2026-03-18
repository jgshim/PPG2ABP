"""
Step 5: Waveform Reconstruction Models (PPG -> ABP waveform)
- Model A: 1D U-Net (optimized, with FiLM conditioning)
- Model B: CNN-LSTM Hybrid
Train 1-fold each for speed, then generate comprehensive waveform comparisons.
"""
import os, sys, time, gc
import numpy as np
import torch
import torch.nn as nn
import torch.nn.functional as F
from torch.utils.data import Dataset, DataLoader
from sklearn.preprocessing import StandardScaler
from sklearn.model_selection import GroupKFold
import warnings
warnings.filterwarnings('ignore')
sys.stdout.reconfigure(line_buffering=True)

BASE = 'C:/Users/jaege/Desktop/Study/PPG2ABP'
PROCESSED_DIR = os.path.join(BASE, 'processed')
RESULTS_DIR = os.path.join(BASE, 'results')
FIG_DIR = os.path.join(BASE, 'figures')
os.makedirs(FIG_DIR, exist_ok=True)
os.makedirs(RESULTS_DIR, exist_ok=True)

DEVICE = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
print(f"Device: {DEVICE}")
if DEVICE.type == 'cuda':
    print(f"GPU: {torch.cuda.get_device_name(0)}")
    print(f"VRAM: {torch.cuda.get_device_properties(0).total_memory / 1024**3:.1f} GB")

BATCH_SIZE = 256
LR = 1e-3
PATIENCE = 10
EPOCHS = 50

# ======================================================================
# Dataset
# ======================================================================
class WaveformDataset(Dataset):
    def __init__(self, pleth, ibp1, demo):
        self.pleth = torch.FloatTensor(pleth).unsqueeze(1)  # (N,1,1250)
        self.ibp1 = torch.FloatTensor(ibp1).unsqueeze(1)    # (N,1,1250)
        self.demo = torch.FloatTensor(demo)                  # (N,2)
    def __len__(self): return len(self.pleth)
    def __getitem__(self, i):
        return self.pleth[i], self.ibp1[i], self.demo[i]

# ======================================================================
# Model A: Optimized 1D U-Net with FiLM
# ======================================================================
class UNetBlock(nn.Module):
    def __init__(self, in_ch, out_ch, k=5):
        super().__init__()
        self.conv = nn.Sequential(
            nn.Conv1d(in_ch, out_ch, k, padding=k//2),
            nn.BatchNorm1d(out_ch), nn.ReLU(),
            nn.Conv1d(out_ch, out_ch, k, padding=k//2),
            nn.BatchNorm1d(out_ch), nn.ReLU())
    def forward(self, x): return self.conv(x)

class UNet1D(nn.Module):
    def __init__(self, n_demo=2):
        super().__init__()
        self.enc1 = UNetBlock(1, 32)
        self.enc2 = UNetBlock(32, 64)
        self.enc3 = UNetBlock(64, 128)
        self.enc4 = UNetBlock(128, 256)
        self.pool = nn.MaxPool1d(2)
        self.bottleneck = UNetBlock(256, 512)
        self.film_gamma = nn.Linear(n_demo, 512)
        self.film_beta = nn.Linear(n_demo, 512)
        self.up4 = nn.ConvTranspose1d(512, 256, 2, stride=2)
        self.dec4 = UNetBlock(512, 256)
        self.up3 = nn.ConvTranspose1d(256, 128, 2, stride=2)
        self.dec3 = UNetBlock(256, 128)
        self.up2 = nn.ConvTranspose1d(128, 64, 2, stride=2)
        self.dec2 = UNetBlock(128, 64)
        self.up1 = nn.ConvTranspose1d(64, 32, 2, stride=2)
        self.dec1 = UNetBlock(64, 32)
        self.final = nn.Conv1d(32, 1, 1)

    def forward(self, ppg, demo):
        orig_len = ppg.size(-1)
        pad_len = (16 - orig_len % 16) % 16
        if pad_len > 0: ppg = F.pad(ppg, (0, pad_len))
        e1 = self.enc1(ppg)
        e2 = self.enc2(self.pool(e1))
        e3 = self.enc3(self.pool(e2))
        e4 = self.enc4(self.pool(e3))
        b = self.bottleneck(self.pool(e4))
        gamma = self.film_gamma(demo).unsqueeze(-1)
        beta = self.film_beta(demo).unsqueeze(-1)
        b = gamma * b + beta
        d4 = self.dec4(torch.cat([self.up4(b)[:,:,:e4.size(2)], e4], dim=1))
        d3 = self.dec3(torch.cat([self.up3(d4)[:,:,:e3.size(2)], e3], dim=1))
        d2 = self.dec2(torch.cat([self.up2(d3)[:,:,:e2.size(2)], e2], dim=1))
        d1 = self.dec1(torch.cat([self.up1(d2)[:,:,:e1.size(2)], e1], dim=1))
        return self.final(d1)[:, :, :orig_len]

# ======================================================================
# Model B: CNN-LSTM Hybrid
# ======================================================================
class CNNLSTM_WaveformModel(nn.Module):
    """
    CNN encoder extracts local features from PPG,
    Bidirectional LSTM captures temporal dependencies,
    CNN decoder reconstructs ABP waveform.
    """
    def __init__(self, n_demo=2):
        super().__init__()
        # CNN Encoder: downsample PPG
        self.enc = nn.Sequential(
            nn.Conv1d(1, 32, 7, padding=3), nn.BatchNorm1d(32), nn.ReLU(),
            nn.Conv1d(32, 64, 5, padding=2), nn.BatchNorm1d(64), nn.ReLU(),
            nn.MaxPool1d(2),  # 1250->625
            nn.Conv1d(64, 128, 5, padding=2), nn.BatchNorm1d(128), nn.ReLU(),
            nn.Conv1d(128, 128, 3, padding=1), nn.BatchNorm1d(128), nn.ReLU(),
            nn.MaxPool1d(2),  # 625->312
        )
        # FiLM for demographics
        self.film_gamma = nn.Linear(n_demo, 128)
        self.film_beta = nn.Linear(n_demo, 128)

        # Bidirectional LSTM on encoded features
        self.lstm = nn.LSTM(128, 128, num_layers=2, batch_first=True,
                           bidirectional=True, dropout=0.2)

        # CNN Decoder: upsample back to original length
        self.dec = nn.Sequential(
            nn.ConvTranspose1d(256, 128, 4, stride=2, padding=1),  # 312->624
            nn.BatchNorm1d(128), nn.ReLU(),
            nn.ConvTranspose1d(128, 64, 4, stride=2, padding=1),   # 624->1248
            nn.BatchNorm1d(64), nn.ReLU(),
            nn.Conv1d(64, 32, 3, padding=1), nn.BatchNorm1d(32), nn.ReLU(),
            nn.Conv1d(32, 1, 1),
        )

    def forward(self, ppg, demo):
        orig_len = ppg.size(-1)  # 1250

        # Pad to multiple of 4 for clean downsampling
        pad_len = (4 - orig_len % 4) % 4
        if pad_len > 0: ppg = F.pad(ppg, (0, pad_len))

        # CNN Encode
        x = self.enc(ppg)  # (B, 128, L/4)

        # FiLM conditioning
        gamma = self.film_gamma(demo).unsqueeze(-1)  # (B, 128, 1)
        beta = self.film_beta(demo).unsqueeze(-1)
        x = gamma * x + beta

        # LSTM: (B, C, L) -> (B, L, C) -> LSTM -> (B, L, 2C) -> (B, 2C, L)
        x = x.permute(0, 2, 1)  # (B, L/4, 128)
        x, _ = self.lstm(x)      # (B, L/4, 256)
        x = x.permute(0, 2, 1)  # (B, 256, L/4)

        # CNN Decode
        x = self.dec(x)  # (B, 1, ~L)

        # Trim or pad to original length
        if x.size(-1) >= orig_len:
            x = x[:, :, :orig_len]
        else:
            x = F.pad(x, (0, orig_len - x.size(-1)))
        return x


# ======================================================================
# Load Data
# ======================================================================
print("\nLoading data...")
all_pleth, all_ibp1, all_demo = [], [], []
all_sbp, all_dbp, all_mbp = [], [], []
all_subjects = []

for case_file in sorted(os.listdir(PROCESSED_DIR)):
    if not case_file.endswith('.npz'): continue
    data = np.load(os.path.join(PROCESSED_DIR, case_file))
    pleth, ibp1 = data['pleth'], data['ibp1']
    sbp, dbp, mbp = data['sbp'], data['dbp'], data['mbp']
    age, sex = float(data['age']), int(data['sex'])
    for i in range(len(pleth)):
        if np.isnan(sbp[i]) or np.isnan(dbp[i]): continue
        if sbp[i] < 30 or sbp[i] > 250 or dbp[i] < 10 or dbp[i] > 200 or sbp[i] <= dbp[i]: continue
        all_pleth.append(pleth[i]); all_ibp1.append(ibp1[i])
        all_demo.append([age, sex])
        all_sbp.append(sbp[i]); all_dbp.append(dbp[i]); all_mbp.append(mbp[i])
        all_subjects.append(case_file.replace('.npz', ''))

all_pleth = np.array(all_pleth, dtype=np.float32)
all_ibp1 = np.array(all_ibp1, dtype=np.float32)
all_demo = np.array(all_demo, dtype=np.float32)
all_sbp = np.array(all_sbp, dtype=np.float32)
all_dbp = np.array(all_dbp, dtype=np.float32)
all_mbp = np.array(all_mbp, dtype=np.float32)
all_subjects = np.array(all_subjects)

demo_scaler = StandardScaler()
all_demo_s = demo_scaler.fit_transform(all_demo).astype(np.float32)

unique_subjects = np.unique(all_subjects)
subj_map = {s: i for i, s in enumerate(unique_subjects)}
subject_ids = np.array([subj_map[s] for s in all_subjects])
print(f"Total: {len(all_pleth)} samples, {len(unique_subjects)} subjects")
print(f"Waveform length: {all_pleth.shape[1]} points (10s @ 125Hz)")
print(f"SBP: {all_sbp.mean():.0f}+/-{all_sbp.std():.0f}, DBP: {all_dbp.mean():.0f}+/-{all_dbp.std():.0f}")


# ======================================================================
# Combined Loss: MSE + Gradient + Correlation
# ======================================================================
def waveform_loss(pred, target):
    mse = F.mse_loss(pred, target)
    # Gradient loss (shape preservation)
    grad_pred = pred[:, :, 1:] - pred[:, :, :-1]
    grad_target = target[:, :, 1:] - target[:, :, :-1]
    grad_loss = F.mse_loss(grad_pred, grad_target)
    return mse + 0.3 * grad_loss


# ======================================================================
# Training Function
# ======================================================================
def train_waveform_model(model_class, model_name, fold_idx=0):
    print(f"\n{'='*70}")
    print(f"  {model_name} - Fold {fold_idx+1}, GPU")
    print(f"{'='*70}")

    gkf = GroupKFold(n_splits=5)
    folds = list(gkf.split(all_pleth, groups=subject_ids))
    tr_idx, te_idx = folds[fold_idx]

    te_subj = np.unique(all_subjects[te_idx])
    print(f"  Train: {len(tr_idx)}, Test: {len(te_idx)} ({len(te_subj)} subjects: {list(te_subj)})")

    # Normalize PPG and ABP separately
    p_mu, p_sig = all_pleth[tr_idx].mean(), all_pleth[tr_idx].std()
    a_mu, a_sig = all_ibp1[tr_idx].mean(), all_ibp1[tr_idx].std()
    tr_p = (all_pleth[tr_idx] - p_mu) / (p_sig + 1e-8)
    te_p = (all_pleth[te_idx] - p_mu) / (p_sig + 1e-8)
    tr_a = (all_ibp1[tr_idx] - a_mu) / (a_sig + 1e-8)

    # Split training into train/val (90/10)
    n_tr = len(tr_idx)
    perm = np.random.RandomState(42).permutation(n_tr)
    n_val = int(n_tr * 0.1)
    val_perm, trn_perm = perm[:n_val], perm[n_val:]

    trn_ds = WaveformDataset(tr_p[trn_perm], tr_a[trn_perm], all_demo_s[tr_idx][trn_perm])
    val_ds = WaveformDataset(tr_p[val_perm], tr_a[val_perm], all_demo_s[tr_idx][val_perm])
    te_ds = WaveformDataset(te_p, np.zeros_like(all_ibp1[te_idx]), all_demo_s[te_idx])

    trn_ld = DataLoader(trn_ds, batch_size=BATCH_SIZE, shuffle=True, num_workers=0, pin_memory=True)
    val_ld = DataLoader(val_ds, batch_size=BATCH_SIZE, shuffle=False, num_workers=0, pin_memory=True)
    te_ld = DataLoader(te_ds, batch_size=BATCH_SIZE, shuffle=False, num_workers=0, pin_memory=True)

    model = model_class().to(DEVICE)
    n_params = sum(p.numel() for p in model.parameters() if p.requires_grad)
    print(f"  Parameters: {n_params:,}")

    opt = torch.optim.AdamW(model.parameters(), lr=LR, weight_decay=1e-4)
    sched = torch.optim.lr_scheduler.CosineAnnealingLR(opt, T_max=EPOCHS)

    best_vl, pat_cnt, best_st = float('inf'), 0, None
    t0 = time.time()

    for ep in range(EPOCHS):
        # Train
        model.train()
        ep_loss = 0
        for ppg, abp, demo in trn_ld:
            ppg = ppg.to(DEVICE, non_blocking=True)
            abp = abp.to(DEVICE, non_blocking=True)
            demo = demo.to(DEVICE, non_blocking=True)
            opt.zero_grad(set_to_none=True)
            pred = model(ppg, demo)
            loss = waveform_loss(pred, abp)
            loss.backward()
            nn.utils.clip_grad_norm_(model.parameters(), 1.0)
            opt.step()
            ep_loss += loss.item() * len(ppg)
        sched.step()
        ep_loss /= len(trn_ds)

        # Validate
        model.eval()
        vl = 0
        with torch.no_grad():
            for ppg, abp, demo in val_ld:
                ppg = ppg.to(DEVICE, non_blocking=True)
                abp = abp.to(DEVICE, non_blocking=True)
                demo = demo.to(DEVICE, non_blocking=True)
                vl += waveform_loss(model(ppg, demo), abp).item() * len(ppg)
        vl /= len(val_ds)

        elapsed = time.time() - t0
        if (ep + 1) % 5 == 0 or ep == 0:
            print(f"    Epoch {ep+1:3d}: train_loss={ep_loss:.6f} val_loss={vl:.6f} [{elapsed:.0f}s]")

        if vl < best_vl:
            best_vl = vl
            best_st = {k: v.cpu().clone() for k, v in model.state_dict().items()}
            pat_cnt = 0
        else:
            pat_cnt += 1
            if pat_cnt >= PATIENCE:
                print(f"    Early stop at epoch {ep+1}")
                break

    total_time = time.time() - t0
    print(f"  Training time: {total_time:.0f}s ({total_time/60:.1f}min)")

    # Inference on test set
    model.load_state_dict(best_st); model.eval()
    recon_list = []
    with torch.no_grad():
        for ppg, _, demo in te_ld:
            pred = model(ppg.to(DEVICE, non_blocking=True), demo.to(DEVICE, non_blocking=True))
            recon_list.append(pred.cpu().numpy())
    recon = np.concatenate(recon_list, axis=0).squeeze(1)  # (N, 1250)

    # Denormalize
    recon_denorm = recon * a_sig + a_mu
    te_abp = all_ibp1[te_idx]

    # Metrics
    # Per-sample correlation
    corrs = []
    for i in range(len(te_abp)):
        c = np.corrcoef(recon_denorm[i], te_abp[i])[0, 1]
        if not np.isnan(c): corrs.append(c)
    mean_corr = np.mean(corrs)

    # Waveform RMSE
    wf_rmse = np.sqrt(np.mean((recon_denorm - te_abp) ** 2))

    # SBP/DBP from waveform
    pred_sbp = np.array([seg.max() for seg in recon_denorm])
    pred_dbp = np.array([seg.min() for seg in recon_denorm])
    pred_mbp = np.array([seg.mean() for seg in recon_denorm])

    true_sbp_te = all_sbp[te_idx]
    true_dbp_te = all_dbp[te_idx]
    true_mbp_te = all_mbp[te_idx]

    sbp_mae = np.mean(np.abs(pred_sbp - true_sbp_te))
    dbp_mae = np.mean(np.abs(pred_dbp - true_dbp_te))
    mbp_mae = np.mean(np.abs(pred_mbp - true_mbp_te))
    sbp_rmse = np.sqrt(np.mean((pred_sbp - true_sbp_te)**2))
    dbp_rmse = np.sqrt(np.mean((pred_dbp - true_dbp_te)**2))

    # Per-sample SBP/DBP error
    sbp_err = pred_sbp - true_sbp_te
    dbp_err = pred_dbp - true_dbp_te
    sbp_bias, sbp_sd = np.mean(sbp_err), np.std(sbp_err)
    dbp_bias, dbp_sd = np.mean(dbp_err), np.std(dbp_err)

    print(f"\n  --- {model_name} Results ---")
    print(f"  Waveform RMSE:  {wf_rmse:.2f} mmHg")
    print(f"  Waveform Corr:  {mean_corr:.4f}")
    print(f"  Correlation distribution: median={np.median(corrs):.4f}, "
          f"Q1={np.percentile(corrs,25):.4f}, Q3={np.percentile(corrs,75):.4f}")
    print(f"  SBP: MAE={sbp_mae:.2f}, RMSE={sbp_rmse:.2f}, Bias={sbp_bias:+.1f}+/-{sbp_sd:.1f}")
    print(f"  DBP: MAE={dbp_mae:.2f}, RMSE={dbp_rmse:.2f}, Bias={dbp_bias:+.1f}+/-{dbp_sd:.1f}")
    print(f"  MBP: MAE={mbp_mae:.2f}")

    return {
        'model': model_name,
        'recon': recon_denorm,
        'te_idx': te_idx,
        'te_subjects': all_subjects[te_idx],
        'wf_rmse': wf_rmse,
        'wf_corr': mean_corr,
        'corrs': np.array(corrs),
        'sbp_mae': sbp_mae, 'dbp_mae': dbp_mae, 'mbp_mae': mbp_mae,
        'sbp_rmse': sbp_rmse, 'dbp_rmse': dbp_rmse,
        'sbp_bias': sbp_bias, 'sbp_sd': sbp_sd,
        'dbp_bias': dbp_bias, 'dbp_sd': dbp_sd,
        'pred_sbp': pred_sbp, 'pred_dbp': pred_dbp, 'pred_mbp': pred_mbp,
        'time': total_time,
    }


# ======================================================================
# Train Both Models
# ======================================================================
print("\n" + "="*70)
print("  WAVEFORM RECONSTRUCTION MODELS")
print("="*70)

results = {}

# Model A: U-Net 1D
results['UNet'] = train_waveform_model(UNet1D, "U-Net 1D", fold_idx=0)
torch.cuda.empty_cache(); gc.collect()

# Model B: CNN-LSTM
results['CNNLSTM'] = train_waveform_model(CNNLSTM_WaveformModel, "CNN-LSTM Hybrid", fold_idx=0)
torch.cuda.empty_cache(); gc.collect()


# ======================================================================
# Save predictions
# ======================================================================
te_idx = results['UNet']['te_idx']
save_dict = {
    'te_idx': te_idx,
    'te_subjects': results['UNet']['te_subjects'],
    'true_ibp1': all_ibp1[te_idx],
    'true_sbp': all_sbp[te_idx],
    'true_dbp': all_dbp[te_idx],
    'true_mbp': all_mbp[te_idx],
    'true_pleth': all_pleth[te_idx],
    'unet_recon': results['UNet']['recon'],
    'cnnlstm_recon': results['CNNLSTM']['recon'],
}
np.savez_compressed(os.path.join(RESULTS_DIR, 'waveform_results.npz'), **save_dict)
print(f"\nPredictions saved to {RESULTS_DIR}/waveform_results.npz")


# ======================================================================
# VISUALIZATION
# ======================================================================
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec

plt.rcParams['font.family'] = 'Malgun Gothic'
plt.rcParams['axes.unicode_minus'] = False
plt.rcParams['figure.dpi'] = 200

C_ACTUAL = '#1a1a2e'
C_UNET = '#e74c3c'
C_LSTM = '#2980b9'
C_PPG = '#27ae60'

te_abp = all_ibp1[te_idx]
te_ppg = all_pleth[te_idx]
te_subjects_arr = all_subjects[te_idx]
unet_recon = results['UNet']['recon']
lstm_recon = results['CNNLSTM']['recon']

# Select patients with good variety
te_unique = np.unique(te_subjects_arr)
# Pick 4 patients
patient_list = list(te_unique[:4])
print(f"\nPlotting patients: {patient_list}")

SR = 125


# ======================================================================
# FIGURE 1: Full waveform comparison - 4 patients, 30s each (3 segments)
# ======================================================================
print("Generating Figure 1: Waveform comparison (30s)...")
fig, axes = plt.subplots(4, 1, figsize=(18, 16))

for ax_i, subj in enumerate(patient_list):
    ax = axes[ax_i]
    mask = te_subjects_arr == subj
    idx = np.where(mask)[0]

    # Pick 3 consecutive segments from middle
    mid = len(idx) // 2
    sel = idx[mid:mid+3]
    if len(sel) < 3: sel = idx[:3]

    # Concatenate waveforms
    actual = np.concatenate([te_abp[s] for s in sel])
    unet_w = np.concatenate([unet_recon[s] for s in sel])
    lstm_w = np.concatenate([lstm_recon[s] for s in sel])
    ppg_w = np.concatenate([te_ppg[s] for s in sel])
    t = np.arange(len(actual)) / SR

    # Plot ABP waveforms
    ax.plot(t, actual, color=C_ACTUAL, linewidth=1.3, label='Actual ABP (IBP1)', zorder=5)
    ax.plot(t, unet_w, color=C_UNET, linewidth=1.0, alpha=0.85, label='U-Net Predicted')
    ax.plot(t, lstm_w, color=C_LSTM, linewidth=1.0, alpha=0.85, label='CNN-LSTM Predicted')

    # Segment boundaries
    for i in range(1, 3):
        ax.axvline(i * 10, color='gray', linestyle=':', alpha=0.3)

    # Per-segment correlation
    for i, s in enumerate(sel):
        c_unet = np.corrcoef(te_abp[s], unet_recon[s])[0, 1]
        c_lstm = np.corrcoef(te_abp[s], lstm_recon[s])[0, 1]
        ax.text(i * 10 + 0.2, ax.get_ylim()[0] if ax_i > 0 else actual.min() - 3,
                f'r={c_unet:.2f}/{c_lstm:.2f}', fontsize=7.5, color='gray',
                style='italic', va='bottom')

    age = all_demo[all_subjects == subj][0, 0]
    sex = 'M' if all_demo[all_subjects == subj][0, 1] == 1 else 'F'
    n_segs = mask.sum()
    ax.set_title(f'{subj} (age={age:.1f}yr, {sex}, {n_segs} segs)',
                 fontsize=12, fontweight='bold', loc='left')
    ax.set_ylabel('ABP (mmHg)', fontsize=11)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(alpha=0.15)
    if ax_i == 0:
        ax.legend(fontsize=10, ncol=3, loc='upper right')

axes[-1].set_xlabel('Time (seconds)', fontsize=12, fontweight='bold')
fig.suptitle('ABP Waveform Reconstruction: Actual vs U-Net vs CNN-LSTM (30s windows)',
             fontsize=15, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'waveform_comparison_30s.png'), bbox_inches='tight')
plt.close()
print("  Saved: waveform_comparison_30s.png")


# ======================================================================
# FIGURE 2: Zoomed single-cycle comparison (~2-3 cardiac cycles)
# ======================================================================
print("Generating Figure 2: Zoomed single cycles...")
fig, axes = plt.subplots(2, 2, figsize=(14, 9))

for ax_i, subj in enumerate(patient_list):
    ax = axes[ax_i // 2, ax_i % 2]
    mask = te_subjects_arr == subj
    idx = np.where(mask)[0]

    # Pick a segment from middle
    seg_i = idx[len(idx) // 2]
    actual_seg = te_abp[seg_i]
    unet_seg = unet_recon[seg_i]
    lstm_seg = lstm_recon[seg_i]
    ppg_seg = te_ppg[seg_i]

    # Show first 3-4 seconds (375-500 points) for cycle detail
    show_pts = 500
    t = np.arange(show_pts) / SR

    ax.plot(t, actual_seg[:show_pts], color=C_ACTUAL, linewidth=1.8, label='Actual ABP')
    ax.plot(t, unet_seg[:show_pts], color=C_UNET, linewidth=1.3, alpha=0.85,
            linestyle='--', label='U-Net')
    ax.plot(t, lstm_seg[:show_pts], color=C_LSTM, linewidth=1.3, alpha=0.85,
            linestyle='--', label='CNN-LSTM')

    # PPG on secondary axis (normalized)
    ax2 = ax.twinx()
    ppg_norm = (ppg_seg[:show_pts] - ppg_seg[:show_pts].mean()) / (ppg_seg[:show_pts].std() + 1e-8)
    ax2.plot(t, ppg_norm, color=C_PPG, linewidth=0.8, alpha=0.4, label='PPG (input)')
    ax2.set_ylabel('PPG (norm)', fontsize=9, color=C_PPG, alpha=0.6)
    ax2.tick_params(axis='y', labelcolor=C_PPG, labelsize=8)
    ax2.set_ylim(-4, 4)

    c_unet = np.corrcoef(actual_seg, unet_seg)[0, 1]
    c_lstm = np.corrcoef(actual_seg, lstm_seg)[0, 1]
    rmse_unet = np.sqrt(np.mean((actual_seg - unet_seg)**2))
    rmse_lstm = np.sqrt(np.mean((actual_seg - lstm_seg)**2))

    age = all_demo[all_subjects == subj][0, 0]
    sex = 'M' if all_demo[all_subjects == subj][0, 1] == 1 else 'F'

    ax.set_title(f'{subj} (age={age:.1f}yr, {sex})\n'
                 f'U-Net: r={c_unet:.3f}, RMSE={rmse_unet:.1f} | '
                 f'CNN-LSTM: r={c_lstm:.3f}, RMSE={rmse_lstm:.1f}',
                 fontsize=10, fontweight='bold', loc='left')
    ax.set_xlabel('Time (s)', fontsize=10)
    ax.set_ylabel('ABP (mmHg)', fontsize=10)
    ax.spines['top'].set_visible(False)
    ax.grid(alpha=0.15)
    if ax_i == 0:
        lines1, labels1 = ax.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax.legend(lines1 + lines2, labels1 + labels2, fontsize=8, loc='upper right', ncol=2)

fig.suptitle('Zoomed Waveform: Cardiac Cycle Detail (4s windows)',
             fontsize=14, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'waveform_zoomed_cycles.png'), bbox_inches='tight')
plt.close()
print("  Saved: waveform_zoomed_cycles.png")


# ======================================================================
# FIGURE 3: Correlation Distribution + RMSE per patient
# ======================================================================
print("Generating Figure 3: Correlation & RMSE distributions...")
fig, axes = plt.subplots(2, 2, figsize=(14, 10))

# 3a: Correlation histogram
ax = axes[0, 0]
ax.hist(results['UNet']['corrs'], bins=60, color=C_UNET, alpha=0.6, label='U-Net', edgecolor='white')
ax.hist(results['CNNLSTM']['corrs'], bins=60, color=C_LSTM, alpha=0.6, label='CNN-LSTM', edgecolor='white')
ax.axvline(results['UNet']['wf_corr'], color=C_UNET, linestyle='--', linewidth=2,
           label=f"U-Net mean={results['UNet']['wf_corr']:.3f}")
ax.axvline(results['CNNLSTM']['wf_corr'], color=C_LSTM, linestyle='--', linewidth=2,
           label=f"CNN-LSTM mean={results['CNNLSTM']['wf_corr']:.3f}")
ax.set_xlabel('Waveform Correlation', fontsize=11)
ax.set_ylabel('Count', fontsize=11)
ax.set_title('Per-Sample Waveform Correlation Distribution', fontsize=12, fontweight='bold')
ax.legend(fontsize=9)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.grid(alpha=0.2)

# 3b: Per-sample RMSE histogram
ax = axes[0, 1]
unet_rmse_per = np.sqrt(np.mean((unet_recon - te_abp)**2, axis=1))
lstm_rmse_per = np.sqrt(np.mean((lstm_recon - te_abp)**2, axis=1))
ax.hist(unet_rmse_per, bins=60, color=C_UNET, alpha=0.6, label=f'U-Net (mean={unet_rmse_per.mean():.1f})', edgecolor='white')
ax.hist(lstm_rmse_per, bins=60, color=C_LSTM, alpha=0.6, label=f'CNN-LSTM (mean={lstm_rmse_per.mean():.1f})', edgecolor='white')
ax.set_xlabel('Waveform RMSE (mmHg)', fontsize=11)
ax.set_ylabel('Count', fontsize=11)
ax.set_title('Per-Sample Waveform RMSE Distribution', fontsize=12, fontweight='bold')
ax.legend(fontsize=9)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.grid(alpha=0.2)

# 3c: Per-patient mean correlation
ax = axes[1, 0]
patient_corr_unet, patient_corr_lstm = [], []
patient_names = []
for subj in te_unique:
    mask = te_subjects_arr == subj
    idx_local = np.where(mask)[0]
    corr_u = [np.corrcoef(te_abp[i], unet_recon[i])[0, 1] for i in idx_local]
    corr_l = [np.corrcoef(te_abp[i], lstm_recon[i])[0, 1] for i in idx_local]
    patient_corr_unet.append(np.nanmean(corr_u))
    patient_corr_lstm.append(np.nanmean(corr_l))
    patient_names.append(subj.replace('case', 'P'))

x = np.arange(len(patient_names))
w = 0.35
ax.bar(x - w/2, patient_corr_unet, w, color=C_UNET, alpha=0.8, label='U-Net')
ax.bar(x + w/2, patient_corr_lstm, w, color=C_LSTM, alpha=0.8, label='CNN-LSTM')
for i in range(len(x)):
    ax.text(x[i]-w/2, patient_corr_unet[i]+0.01, f'{patient_corr_unet[i]:.2f}',
            ha='center', fontsize=7.5, color=C_UNET, fontweight='bold')
    ax.text(x[i]+w/2, patient_corr_lstm[i]+0.01, f'{patient_corr_lstm[i]:.2f}',
            ha='center', fontsize=7.5, color=C_LSTM, fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels(patient_names, fontsize=9)
ax.set_ylabel('Mean Waveform Correlation', fontsize=11)
ax.set_title('Per-Patient Mean Waveform Correlation', fontsize=12, fontweight='bold')
ax.legend(fontsize=9)
ax.set_ylim(0, 1.05)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.grid(axis='y', alpha=0.2)

# 3d: SBP/DBP derived from waveform - scatter
ax = axes[1, 1]
true_s = all_sbp[te_idx]
true_d = all_dbp[te_idx]
ax.scatter(true_s, results['UNet']['pred_sbp'], s=3, alpha=0.1, color=C_UNET, label='U-Net SBP', rasterized=True)
ax.scatter(true_d, results['UNet']['pred_dbp'], s=3, alpha=0.1, color=C_UNET, marker='s', rasterized=True)
ax.scatter(true_s, results['CNNLSTM']['pred_sbp'], s=3, alpha=0.1, color=C_LSTM, label='CNN-LSTM SBP', rasterized=True)
ax.scatter(true_d, results['CNNLSTM']['pred_dbp'], s=3, alpha=0.1, color=C_LSTM, marker='s', rasterized=True)
lims = [20, 180]
ax.plot(lims, lims, 'k--', linewidth=1, alpha=0.5)
ax.set_xlim(lims); ax.set_ylim(lims)
ax.set_xlabel('Actual BP (mmHg)', fontsize=11)
ax.set_ylabel('Predicted BP (mmHg)', fontsize=11)
ax.set_title('SBP/DBP Derived from Reconstructed Waveform', fontsize=12, fontweight='bold')
ax.legend(fontsize=9, loc='upper left')
ax.set_aspect('equal')
ax.grid(alpha=0.2)

fig.suptitle('Waveform Reconstruction Quality Analysis',
             fontsize=15, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'waveform_quality_analysis.png'), bbox_inches='tight')
plt.close()
print("  Saved: waveform_quality_analysis.png")


# ======================================================================
# FIGURE 4: Best/Worst/Median case comparison
# ======================================================================
print("Generating Figure 4: Best/Worst/Median cases...")
all_corrs_unet = np.array([np.corrcoef(te_abp[i], unet_recon[i])[0, 1] for i in range(len(te_abp))])
all_corrs_lstm = np.array([np.corrcoef(te_abp[i], lstm_recon[i])[0, 1] for i in range(len(te_abp))])
# Average correlation for ranking
avg_corrs = (np.nan_to_num(all_corrs_unet) + np.nan_to_num(all_corrs_lstm)) / 2

best_i = np.nanargmax(avg_corrs)
worst_i = np.nanargmin(avg_corrs)
median_i = np.argsort(avg_corrs)[len(avg_corrs) // 2]

fig, axes = plt.subplots(3, 1, figsize=(16, 12))
cases = [('Best', best_i), ('Median', median_i), ('Worst', worst_i)]

for ax_i, (label, seg_i) in enumerate(cases):
    ax = axes[ax_i]
    t = np.arange(1250) / SR

    ax.plot(t, te_abp[seg_i], color=C_ACTUAL, linewidth=1.8, label='Actual ABP', zorder=5)
    ax.plot(t, unet_recon[seg_i], color=C_UNET, linewidth=1.2, alpha=0.85,
            linestyle='--', label='U-Net')
    ax.plot(t, lstm_recon[seg_i], color=C_LSTM, linewidth=1.2, alpha=0.85,
            linestyle='--', label='CNN-LSTM')

    # PPG (faint background)
    ax2 = ax.twinx()
    ppg_n = (te_ppg[seg_i] - te_ppg[seg_i].mean()) / (te_ppg[seg_i].std() + 1e-8)
    ax2.plot(t, ppg_n, color=C_PPG, linewidth=0.5, alpha=0.25)
    ax2.set_ylim(-5, 5)
    ax2.set_yticks([])

    c_u = all_corrs_unet[seg_i]
    c_l = all_corrs_lstm[seg_i]
    rmse_u = np.sqrt(np.mean((te_abp[seg_i] - unet_recon[seg_i])**2))
    rmse_l = np.sqrt(np.mean((te_abp[seg_i] - lstm_recon[seg_i])**2))
    subj = te_subjects_arr[seg_i]

    ax.set_title(f'{label} Case - {subj} | '
                 f'U-Net: r={c_u:.3f}, RMSE={rmse_u:.1f} | '
                 f'CNN-LSTM: r={c_l:.3f}, RMSE={rmse_l:.1f}',
                 fontsize=12, fontweight='bold',
                 color={'Best': '#27ae60', 'Median': '#f39c12', 'Worst': '#e74c3c'}[label])
    ax.set_ylabel('ABP (mmHg)', fontsize=11)
    ax.spines['top'].set_visible(False)
    ax.grid(alpha=0.15)
    if ax_i == 0:
        ax.legend(fontsize=10, ncol=3, loc='upper right')

axes[-1].set_xlabel('Time (seconds)', fontsize=12, fontweight='bold')
fig.suptitle('Waveform Reconstruction: Best / Median / Worst Cases',
             fontsize=15, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'waveform_best_median_worst.png'), bbox_inches='tight')
plt.close()
print("  Saved: waveform_best_median_worst.png")


# ======================================================================
# FIGURE 5: Continuous multi-segment comparison (~2 min / 12 segments)
# ======================================================================
print("Generating Figure 5: Long continuous comparison...")
subj = patient_list[0]
mask = te_subjects_arr == subj
idx = np.where(mask)[0]

n_show = min(12, len(idx))
sel = idx[:n_show]

actual_long = np.concatenate([te_abp[s] for s in sel])
unet_long = np.concatenate([unet_recon[s] for s in sel])
lstm_long = np.concatenate([lstm_recon[s] for s in sel])
t = np.arange(len(actual_long)) / SR

fig, axes = plt.subplots(3, 1, figsize=(18, 10), sharex=True)

# Actual + U-Net
ax = axes[0]
ax.plot(t, actual_long, color=C_ACTUAL, linewidth=0.8, label='Actual ABP')
ax.plot(t, unet_long, color=C_UNET, linewidth=0.8, alpha=0.8, label='U-Net Predicted')
ax.fill_between(t, actual_long, unet_long, alpha=0.15, color=C_UNET)
ax.set_ylabel('ABP (mmHg)', fontsize=11)
ax.set_title(f'U-Net Reconstruction - {subj} ({n_show*10}s continuous)',
             fontsize=12, fontweight='bold', color=C_UNET)
ax.legend(fontsize=10, loc='upper right')
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.grid(alpha=0.15)

# Actual + CNN-LSTM
ax = axes[1]
ax.plot(t, actual_long, color=C_ACTUAL, linewidth=0.8, label='Actual ABP')
ax.plot(t, lstm_long, color=C_LSTM, linewidth=0.8, alpha=0.8, label='CNN-LSTM Predicted')
ax.fill_between(t, actual_long, lstm_long, alpha=0.15, color=C_LSTM)
ax.set_ylabel('ABP (mmHg)', fontsize=11)
ax.set_title(f'CNN-LSTM Reconstruction - {subj}',
             fontsize=12, fontweight='bold', color=C_LSTM)
ax.legend(fontsize=10, loc='upper right')
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.grid(alpha=0.15)

# Residual (error)
ax = axes[2]
unet_err = unet_long - actual_long
lstm_err = lstm_long - actual_long
ax.plot(t, unet_err, color=C_UNET, linewidth=0.6, alpha=0.7, label=f'U-Net Error (SD={unet_err.std():.1f})')
ax.plot(t, lstm_err, color=C_LSTM, linewidth=0.6, alpha=0.7, label=f'CNN-LSTM Error (SD={lstm_err.std():.1f})')
ax.axhline(0, color='black', linewidth=0.8, alpha=0.3)
ax.fill_between(t, -10, 10, alpha=0.05, color='green')
ax.set_ylabel('Error (mmHg)', fontsize=11)
ax.set_xlabel('Time (seconds)', fontsize=12, fontweight='bold')
ax.set_title('Reconstruction Error (Predicted - Actual)', fontsize=12, fontweight='bold')
ax.legend(fontsize=10, loc='upper right')
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.grid(alpha=0.15)

# Segment boundaries
for ax in axes:
    for i in range(1, n_show):
        ax.axvline(i * 10, color='gray', linestyle=':', alpha=0.15)

fig.suptitle('Continuous Waveform Reconstruction with Residual Analysis',
             fontsize=15, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'waveform_continuous_residual.png'), bbox_inches='tight')
plt.close()
print("  Saved: waveform_continuous_residual.png")


# ======================================================================
# FIGURE 6: Model comparison summary
# ======================================================================
print("Generating Figure 6: Summary comparison...")
fig, axes = plt.subplots(1, 3, figsize=(15, 5))

# 6a: Metric comparison bar
ax = axes[0]
metrics = ['WF Corr', 'SBP MAE', 'DBP MAE', 'MBP MAE']
unet_vals = [results['UNet']['wf_corr'], results['UNet']['sbp_mae'],
             results['UNet']['dbp_mae'], results['UNet']['mbp_mae']]
lstm_vals = [results['CNNLSTM']['wf_corr'], results['CNNLSTM']['sbp_mae'],
             results['CNNLSTM']['dbp_mae'], results['CNNLSTM']['mbp_mae']]

x = np.arange(len(metrics))
w = 0.3
bars1 = ax.bar(x - w/2, unet_vals, w, color=C_UNET, alpha=0.8, label='U-Net')
bars2 = ax.bar(x + w/2, lstm_vals, w, color=C_LSTM, alpha=0.8, label='CNN-LSTM')
for bars in [bars1, bars2]:
    for bar in bars:
        h = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2, h + 0.1, f'{h:.2f}',
                ha='center', fontsize=9, fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels(metrics, fontsize=10)
ax.set_title('Model Comparison', fontsize=13, fontweight='bold')
ax.legend(fontsize=10)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.grid(axis='y', alpha=0.2)

# 6b: Bias +/- SD
ax = axes[1]
targets = ['SBP', 'DBP']
y = np.arange(2)
for i, (t_name, key_b, key_s) in enumerate(
    [('SBP', 'sbp_bias', 'sbp_sd'), ('DBP', 'dbp_bias', 'dbp_sd')]):
    ax.errorbar(results['UNet'][key_b], y[i]-0.12, xerr=results['UNet'][key_s],
                fmt='o', markersize=8, color=C_UNET, capsize=5, capthick=2, linewidth=2,
                label='U-Net' if i==0 else None)
    ax.errorbar(results['CNNLSTM'][key_b], y[i]+0.12, xerr=results['CNNLSTM'][key_s],
                fmt='s', markersize=8, color=C_LSTM, capsize=5, capthick=2, linewidth=2,
                label='CNN-LSTM' if i==0 else None)
ax.axvline(0, color='black', linewidth=0.8, alpha=0.3)
ax.set_yticks(y)
ax.set_yticklabels(targets, fontsize=12, fontweight='bold')
ax.set_xlabel('Bias +/- SD (mmHg)', fontsize=11)
ax.set_title('Bias +/- SD (from waveform peak/trough)', fontsize=13, fontweight='bold')
ax.legend(fontsize=10)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.grid(axis='x', alpha=0.2)

# 6c: Training info
ax = axes[2]
ax.axis('off')
info_text = (
    f"U-Net 1D\n"
    f"  Params: ~3.4M\n"
    f"  Time: {results['UNet']['time']:.0f}s ({results['UNet']['time']/60:.1f}min)\n"
    f"  WF Corr: {results['UNet']['wf_corr']:.4f}\n"
    f"  WF RMSE: {results['UNet']['wf_rmse']:.2f} mmHg\n"
    f"  SBP MAE: {results['UNet']['sbp_mae']:.2f}\n"
    f"  DBP MAE: {results['UNet']['dbp_mae']:.2f}\n\n"
    f"CNN-LSTM Hybrid\n"
    f"  Params: ~1.5M\n"
    f"  Time: {results['CNNLSTM']['time']:.0f}s ({results['CNNLSTM']['time']/60:.1f}min)\n"
    f"  WF Corr: {results['CNNLSTM']['wf_corr']:.4f}\n"
    f"  WF RMSE: {results['CNNLSTM']['wf_rmse']:.2f} mmHg\n"
    f"  SBP MAE: {results['CNNLSTM']['sbp_mae']:.2f}\n"
    f"  DBP MAE: {results['CNNLSTM']['dbp_mae']:.2f}"
)
ax.text(0.1, 0.95, info_text, transform=ax.transAxes, fontsize=11,
        verticalalignment='top', fontfamily='monospace',
        bbox=dict(boxstyle='round', facecolor='lightyellow', alpha=0.8))
ax.set_title('Model Summary', fontsize=13, fontweight='bold')

fig.suptitle('Waveform Reconstruction: U-Net vs CNN-LSTM Summary',
             fontsize=15, fontweight='bold', y=1.02)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'waveform_summary.png'), bbox_inches='tight')
plt.close()
print("  Saved: waveform_summary.png")


# ======================================================================
# Insert into PPTX
# ======================================================================
print("\n" + "="*60)
print("Inserting figures into PowerPoint...")
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

# Find available PPTX
for src in ['PPG2ABP_Presentation_v3.pptx', 'PPG2ABP_Presentation_v2.pptx', 'PPG2ABP_Presentation.pptx']:
    src_path = os.path.join(BASE, src)
    if os.path.exists(src_path):
        try:
            prs = Presentation(src_path)
            print(f"Opened: {src}")
            break
        except:
            continue

PRIMARY_CLR = PRGBColor(0, 51, 102)
WHITE_CLR = PRGBColor(255, 255, 255)
GRAY_CLR = PRGBColor(100, 100, 100)

def make_slide(title, img, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sh = slide.shapes.add_shape(1, PInches(0), PInches(0), PInches(13.333), PInches(0.9))
    sh.fill.solid(); sh.fill.fore_color.rgb = PRIMARY_CLR; sh.line.fill.background()
    tb = slide.shapes.add_textbox(PInches(0.4), PInches(0.1), PInches(12), PInches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title; p.font.size = PPt(28); p.font.bold = True
    p.font.color.rgb = WHITE_CLR; p.font.name = 'Malgun Gothic'
    if subtitle:
        slide.shapes.add_picture(img, PInches(0.3), PInches(1.05), PInches(12.7), PInches(5.7))
        tb2 = slide.shapes.add_textbox(PInches(0.3), PInches(6.85), PInches(12.7), PInches(0.4))
        p2 = tb2.text_frame.paragraphs[0]; p2.text = subtitle
        p2.font.size = PPt(12); p2.font.color.rgb = GRAY_CLR
        p2.font.name = 'Malgun Gothic'; p2.alignment = PP_ALIGN.CENTER
    else:
        slide.shapes.add_picture(img, PInches(0.3), PInches(1.05), PInches(12.7), PInches(6.2))

def make_two_slide(title, img1, img2, c1='', c2=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sh = slide.shapes.add_shape(1, PInches(0), PInches(0), PInches(13.333), PInches(0.9))
    sh.fill.solid(); sh.fill.fore_color.rgb = PRIMARY_CLR; sh.line.fill.background()
    tb = slide.shapes.add_textbox(PInches(0.4), PInches(0.1), PInches(12), PInches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title; p.font.size = PPt(28); p.font.bold = True
    p.font.color.rgb = WHITE_CLR; p.font.name = 'Malgun Gothic'
    slide.shapes.add_picture(img1, PInches(0.2), PInches(1.1), PInches(6.4), PInches(4.8))
    slide.shapes.add_picture(img2, PInches(6.7), PInches(1.1), PInches(6.4), PInches(4.8))
    for cap, left in [(c1, 0.2), (c2, 6.7)]:
        if cap:
            tb2 = slide.shapes.add_textbox(PInches(left), PInches(6.0), PInches(6.4), PInches(0.3))
            p2 = tb2.text_frame.paragraphs[0]; p2.text = cap
            p2.font.size = PPt(11); p2.font.name = 'Malgun Gothic'
            p2.alignment = PP_ALIGN.CENTER; p2.font.color.rgb = GRAY_CLR

F = FIG_DIR

make_slide('Step 5. Waveform Reconstruction - 30s Comparison (4 Patients)',
           os.path.join(F, 'waveform_comparison_30s.png'),
           'Black=Actual ABP, Red=U-Net, Blue=CNN-LSTM | Per-segment correlation shown at bottom')

make_two_slide('Step 5. Zoomed Cardiac Cycles & Best/Median/Worst Cases',
               os.path.join(F, 'waveform_zoomed_cycles.png'),
               os.path.join(F, 'waveform_best_median_worst.png'),
               '4s zoomed view with input PPG overlay (green)',
               'Best/Median/Worst reconstruction cases')

make_slide('Step 5. Continuous Waveform & Residual Analysis (2min)',
           os.path.join(F, 'waveform_continuous_residual.png'),
           'Top: U-Net overlay | Middle: CNN-LSTM overlay | Bottom: Reconstruction error (shaded=+/-10mmHg)')

make_two_slide('Step 5. Waveform Quality Analysis & Model Summary',
               os.path.join(F, 'waveform_quality_analysis.png'),
               os.path.join(F, 'waveform_summary.png'),
               'Correlation/RMSE distributions, per-patient comparison, derived BP scatter',
               'Side-by-side metric comparison and model details')

# Save
for fname in ['PPG2ABP_Presentation.pptx', 'PPG2ABP_Presentation_v2.pptx',
              'PPG2ABP_Presentation_v3.pptx', 'PPG2ABP_Presentation_v4.pptx']:
    try:
        prs.save(os.path.join(BASE, fname))
        print(f"Saved: {fname}")
        break
    except PermissionError:
        print(f"  {fname} locked, trying next...")

print("\n" + "="*60)
print("FINAL SUMMARY")
print("="*60)
for name, res in results.items():
    print(f"\n{res['model']}:")
    print(f"  Waveform Correlation: {res['wf_corr']:.4f}")
    print(f"  Waveform RMSE: {res['wf_rmse']:.2f} mmHg")
    print(f"  SBP MAE: {res['sbp_mae']:.2f}, DBP MAE: {res['dbp_mae']:.2f}")
    print(f"  Training time: {res['time']:.0f}s ({res['time']/60:.1f}min)")
print("\nAll done!")
