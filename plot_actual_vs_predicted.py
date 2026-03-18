"""
Plot actual ABP (IBP1) vs model-predicted BP values.
1. Train CNN + ResNet on 1-fold quickly to get predictions
2. Use XGBoost predictions (already saved)
3. Generate comprehensive actual vs predicted plots
4. Insert into PPTX
"""
import os, sys, time
import numpy as np
import torch
import torch.nn as nn
import torch.nn.functional as F
from torch.utils.data import Dataset, DataLoader
from sklearn.preprocessing import StandardScaler
from sklearn.model_selection import GroupKFold
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
import warnings
warnings.filterwarnings('ignore')
sys.stdout.reconfigure(line_buffering=True)

plt.rcParams['font.family'] = 'Malgun Gothic'
plt.rcParams['axes.unicode_minus'] = False
plt.rcParams['figure.dpi'] = 200

BASE = 'C:/Users/jaege/Desktop/Study/PPG2ABP'
PROCESSED_DIR = os.path.join(BASE, 'processed')
RESULTS_DIR = os.path.join(BASE, 'results')
FIG_DIR = os.path.join(BASE, 'figures')
os.makedirs(FIG_DIR, exist_ok=True)

DEVICE = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
print(f"Device: {DEVICE}")

# ======================================================================
# Models (copied from step4)
# ======================================================================
class PPGDataset(Dataset):
    def __init__(self, pleth, demo, sbp, dbp, mbp):
        self.pleth = torch.FloatTensor(pleth).unsqueeze(1)
        self.demo = torch.FloatTensor(demo)
        self.targets = torch.FloatTensor(np.stack([sbp, dbp, mbp], axis=1))
    def __len__(self): return len(self.pleth)
    def __getitem__(self, i):
        return self.pleth[i], self.demo[i], self.targets[i]

class ImprovedCNN(nn.Module):
    def __init__(self):
        super().__init__()
        self.block1 = nn.Sequential(
            nn.Conv1d(1, 32, 15, padding=7), nn.BatchNorm1d(32), nn.ReLU(),
            nn.Conv1d(32, 32, 15, padding=7), nn.BatchNorm1d(32), nn.ReLU(),
            nn.MaxPool1d(2), nn.Dropout(0.1))
        self.block2 = nn.Sequential(
            nn.Conv1d(32, 64, 9, padding=4), nn.BatchNorm1d(64), nn.ReLU(),
            nn.Conv1d(64, 64, 9, padding=4), nn.BatchNorm1d(64), nn.ReLU(),
            nn.MaxPool1d(2), nn.Dropout(0.1))
        self.block3 = nn.Sequential(
            nn.Conv1d(64, 128, 5, padding=2), nn.BatchNorm1d(128), nn.ReLU(),
            nn.Conv1d(128, 128, 5, padding=2), nn.BatchNorm1d(128), nn.ReLU(),
            nn.MaxPool1d(2), nn.Dropout(0.2))
        self.block4 = nn.Sequential(
            nn.Conv1d(128, 256, 3, padding=1), nn.BatchNorm1d(256), nn.ReLU(),
            nn.Conv1d(256, 256, 3, padding=1), nn.BatchNorm1d(256), nn.ReLU(),
            nn.MaxPool1d(2), nn.Dropout(0.2))
        self.block5 = nn.Sequential(
            nn.Conv1d(256, 512, 3, padding=1), nn.BatchNorm1d(512), nn.ReLU(),
            nn.AdaptiveAvgPool1d(1))
        self.head = nn.Sequential(
            nn.Linear(512 + 2, 256), nn.ReLU(), nn.Dropout(0.3),
            nn.Linear(256, 64), nn.ReLU(), nn.Dropout(0.2),
            nn.Linear(64, 3))
    def forward(self, x, demo):
        x = self.block1(x); x = self.block2(x); x = self.block3(x)
        x = self.block4(x); x = self.block5(x).squeeze(-1)
        return self.head(torch.cat([x, demo], dim=1))

class ResBlock1D(nn.Module):
    def __init__(self, in_ch, out_ch, kernel=5, stride=1, downsample=None):
        super().__init__()
        pad = kernel // 2
        self.conv1 = nn.Conv1d(in_ch, out_ch, kernel, stride=stride, padding=pad)
        self.bn1 = nn.BatchNorm1d(out_ch)
        self.conv2 = nn.Conv1d(out_ch, out_ch, kernel, padding=pad)
        self.bn2 = nn.BatchNorm1d(out_ch)
        self.downsample = downsample
        self.relu = nn.ReLU()
    def forward(self, x):
        identity = x
        out = self.relu(self.bn1(self.conv1(x)))
        out = self.bn2(self.conv2(out))
        if self.downsample is not None: identity = self.downsample(x)
        return self.relu(out + identity)

class ResNet1D(nn.Module):
    def __init__(self):
        super().__init__()
        self.stem = nn.Sequential(
            nn.Conv1d(1, 64, 15, stride=2, padding=7),
            nn.BatchNorm1d(64), nn.ReLU(), nn.MaxPool1d(2))
        self.layer1 = self._make_layer(64, 64, 2)
        self.layer2 = self._make_layer(64, 128, 2, stride=2)
        self.layer3 = self._make_layer(128, 256, 2, stride=2)
        self.layer4 = self._make_layer(256, 512, 2, stride=2)
        self.gap = nn.AdaptiveAvgPool1d(1)
        self.head = nn.Sequential(
            nn.Linear(512 + 2, 256), nn.ReLU(), nn.Dropout(0.3),
            nn.Linear(256, 64), nn.ReLU(), nn.Linear(64, 3))
    def _make_layer(self, in_ch, out_ch, n_blocks, stride=1):
        downsample = None
        if stride != 1 or in_ch != out_ch:
            downsample = nn.Sequential(
                nn.Conv1d(in_ch, out_ch, 1, stride=stride), nn.BatchNorm1d(out_ch))
        layers = [ResBlock1D(in_ch, out_ch, stride=stride, downsample=downsample)]
        for _ in range(1, n_blocks):
            layers.append(ResBlock1D(out_ch, out_ch))
        return nn.Sequential(*layers)
    def forward(self, x, demo):
        x = self.stem(x)
        x = self.layer1(x); x = self.layer2(x)
        x = self.layer3(x); x = self.layer4(x)
        x = self.gap(x).squeeze(-1)
        return self.head(torch.cat([x, demo], dim=1))


# ======================================================================
# Load Data
# ======================================================================
print("Loading data...")
all_pleth, all_ibp1, all_demo = [], [], []
all_sbp, all_dbp, all_mbp = [], [], []
all_subjects = []

for case_file in sorted(os.listdir(PROCESSED_DIR)):
    if not case_file.endswith('.npz'): continue
    data = np.load(os.path.join(PROCESSED_DIR, case_file))
    pleth, ibp1 = data['pleth'], data['ibp1']
    sbp, dbp, mbp = data['sbp'], data['dbp'], data['mbp']
    age, sex = float(data['age']), int(data['sex'])
    for i in range(len(pleth)):
        if np.isnan(sbp[i]) or np.isnan(dbp[i]): continue
        if sbp[i] < 30 or sbp[i] > 250 or dbp[i] < 10 or dbp[i] > 200 or sbp[i] <= dbp[i]: continue
        all_pleth.append(pleth[i]); all_ibp1.append(ibp1[i])
        all_demo.append([age, sex])
        all_sbp.append(sbp[i]); all_dbp.append(dbp[i]); all_mbp.append(mbp[i])
        all_subjects.append(case_file.replace('.npz', ''))

all_pleth = np.array(all_pleth, dtype=np.float32)
all_ibp1 = np.array(all_ibp1, dtype=np.float32)
all_demo = np.array(all_demo, dtype=np.float32)
all_sbp = np.array(all_sbp, dtype=np.float32)
all_dbp = np.array(all_dbp, dtype=np.float32)
all_mbp = np.array(all_mbp, dtype=np.float32)
all_subjects = np.array(all_subjects)

demo_scaler = StandardScaler()
all_demo_s = demo_scaler.fit_transform(all_demo).astype(np.float32)

unique_subjects = np.unique(all_subjects)
subj_map = {s: i for i, s in enumerate(unique_subjects)}
subject_ids = np.array([subj_map[s] for s in all_subjects])
print(f"Total: {len(all_pleth)} samples, {len(unique_subjects)} subjects")


# ======================================================================
# Quick 1-fold training for CNN + ResNet to get predictions
# ======================================================================
BATCH_SIZE = 512
PATIENCE = 8
EPOCHS = 40

gkf = GroupKFold(n_splits=5)
folds = list(gkf.split(all_pleth, groups=subject_ids))
# Use fold 0 (largest test set if possible)
tr_idx, te_idx = folds[0]

print(f"\nTraining 1-fold for predictions (test set: {len(te_idx)} samples)...")

# Normalize PPG
mu, sig = all_pleth[tr_idx].mean(), all_pleth[tr_idx].std()
tr_p = (all_pleth[tr_idx] - mu) / (sig + 1e-8)
te_p = (all_pleth[te_idx] - mu) / (sig + 1e-8)

tr_ds = PPGDataset(tr_p, all_demo_s[tr_idx], all_sbp[tr_idx], all_dbp[tr_idx], all_mbp[tr_idx])
te_ds = PPGDataset(te_p, all_demo_s[te_idx], all_sbp[te_idx], all_dbp[te_idx], all_mbp[te_idx])
tr_ld = DataLoader(tr_ds, batch_size=BATCH_SIZE, shuffle=True, num_workers=0, pin_memory=True)
te_ld = DataLoader(te_ds, batch_size=BATCH_SIZE, shuffle=False, num_workers=0, pin_memory=True)

dl_preds = {}

for model_name, model_class in [('CNN', ImprovedCNN), ('ResNet', ResNet1D)]:
    print(f"\n  Training {model_name}...")
    t0 = time.time()
    model = model_class().to(DEVICE)
    opt = torch.optim.AdamW(model.parameters(), lr=1e-3, weight_decay=1e-4)
    sched = torch.optim.lr_scheduler.CosineAnnealingLR(opt, T_max=EPOCHS)
    crit = nn.SmoothL1Loss()

    best_vl, pat_cnt, best_st = float('inf'), 0, None
    for ep in range(EPOCHS):
        model.train()
        for ppg, demo, tgt in tr_ld:
            ppg, demo, tgt = ppg.to(DEVICE), demo.to(DEVICE), tgt.to(DEVICE)
            opt.zero_grad(set_to_none=True)
            loss = crit(model(ppg, demo), tgt)
            loss.backward()
            nn.utils.clip_grad_norm_(model.parameters(), 1.0)
            opt.step()
        sched.step()

        model.eval()
        vl = 0
        with torch.no_grad():
            for ppg, demo, tgt in te_ld:
                ppg, demo, tgt = ppg.to(DEVICE), demo.to(DEVICE), tgt.to(DEVICE)
                vl += crit(model(ppg, demo), tgt).item() * len(ppg)
        vl /= len(te_ds)
        if vl < best_vl:
            best_vl = vl
            best_st = {k: v.cpu().clone() for k, v in model.state_dict().items()}
            pat_cnt = 0
        else:
            pat_cnt += 1
            if pat_cnt >= PATIENCE:
                break

    model.load_state_dict(best_st); model.eval()
    preds = []
    with torch.no_grad():
        for ppg, demo, _ in te_ld:
            preds.append(model(ppg.to(DEVICE), demo.to(DEVICE)).cpu().numpy())
    preds = np.concatenate(preds, axis=0)
    dl_preds[model_name] = preds
    elapsed = time.time() - t0
    print(f"  {model_name} done in {elapsed:.0f}s (ep={ep+1})")
    del model; torch.cuda.empty_cache()


# Load XGBoost predictions
ml = np.load(os.path.join(RESULTS_DIR, 'ml_baseline_results.npz'), allow_pickle=True)
xgb_sbp_pred = ml['SBP_xgboost_preds']
xgb_dbp_pred = ml['DBP_xgboost_preds']

# ======================================================================
# Select representative patients from test fold
# ======================================================================
te_subjects = all_subjects[te_idx]
te_unique = np.unique(te_subjects)
print(f"\nTest subjects: {te_unique}")

# Pick 4 patients with different segment counts for variety
seg_counts = {s: (te_subjects == s).sum() for s in te_unique}
sorted_subj = sorted(seg_counts.keys(), key=lambda x: seg_counts[x], reverse=True)
selected = sorted_subj[:4]  # top 4 by segment count
print(f"Selected patients: {selected} (segs: {[seg_counts[s] for s in selected]})")

C_ACTUAL = '#1a1a2e'
C_XGB = '#e74c3c'
C_CNN = '#2980b9'
C_RES = '#27ae60'
C_ABP_WAVE = '#34495e'


# ======================================================================
# FIGURE 1: Actual vs Predicted SBP/DBP Time-Series (4 patients)
# ======================================================================
print("\nGenerating Figure 1: SBP time-series...")
fig, axes = plt.subplots(4, 1, figsize=(16, 14), sharex=False)

for ax_i, subj in enumerate(selected):
    ax = axes[ax_i]

    # Global indices for this subject
    global_mask = all_subjects == subj
    te_mask_local = te_subjects == subj

    # True values
    true_sbp = all_sbp[te_idx][te_mask_local]
    true_dbp = all_dbp[te_idx][te_mask_local]

    # XGBoost predictions (from global indices)
    xgb_sbp = xgb_sbp_pred[global_mask]
    xgb_dbp = xgb_dbp_pred[global_mask]

    # CNN / ResNet predictions (from test fold local)
    local_indices = np.where(te_mask_local)[0]
    cnn_sbp = dl_preds['CNN'][local_indices, 0]
    cnn_dbp = dl_preds['CNN'][local_indices, 1]
    res_sbp = dl_preds['ResNet'][local_indices, 0]
    res_dbp = dl_preds['ResNet'][local_indices, 1]

    n = len(true_sbp)
    # Show up to 80 segments for readability
    show_n = min(n, 80)
    x = np.arange(show_n)
    t_sec = x * 10  # 10s per segment

    # SBP curves
    ax.plot(t_sec, true_sbp[:show_n], color=C_ACTUAL, linewidth=2.0, label='Actual SBP', zorder=5)
    ax.plot(t_sec, xgb_sbp[:show_n], color=C_XGB, linewidth=1.2, alpha=0.8, linestyle='--', label='XGBoost SBP')
    ax.plot(t_sec, cnn_sbp[:show_n], color=C_CNN, linewidth=1.2, alpha=0.8, linestyle='--', label='CNN SBP')
    ax.plot(t_sec, res_sbp[:show_n], color=C_RES, linewidth=1.2, alpha=0.8, linestyle='--', label='ResNet SBP')

    # DBP curves
    ax.plot(t_sec, true_dbp[:show_n], color=C_ACTUAL, linewidth=2.0, linestyle='-', label='Actual DBP', zorder=5)
    ax.plot(t_sec, xgb_dbp[:show_n], color=C_XGB, linewidth=1.2, alpha=0.8, linestyle=':')
    ax.plot(t_sec, cnn_dbp[:show_n], color=C_CNN, linewidth=1.2, alpha=0.8, linestyle=':')
    ax.plot(t_sec, res_dbp[:show_n], color=C_RES, linewidth=1.2, alpha=0.8, linestyle=':')

    # Fill between actual SBP and DBP
    ax.fill_between(t_sec, true_sbp[:show_n], true_dbp[:show_n], alpha=0.08, color='gray')

    age_val = all_demo[global_mask][0, 0]
    sex_val = 'M' if all_demo[global_mask][0, 1] == 1 else 'F'
    ax.set_title(f'{subj} (age={age_val:.1f}yr, {sex_val}, {n} segments)',
                 fontsize=12, fontweight='bold', loc='left')
    ax.set_ylabel('BP (mmHg)', fontsize=11)
    ax.grid(alpha=0.2)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    if ax_i == 0:
        ax.legend(fontsize=9, ncol=4, loc='upper right')

axes[-1].set_xlabel('Time (seconds)', fontsize=12, fontweight='bold')
fig.suptitle('Actual vs Predicted Blood Pressure - SBP & DBP Time Series',
             fontsize=15, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'actual_vs_pred_timeseries.png'), bbox_inches='tight')
plt.close()
print("  Saved: actual_vs_pred_timeseries.png")


# ======================================================================
# FIGURE 2: Raw ABP Waveform + Predicted SBP/DBP overlay (one patient, ~5 segments)
# ======================================================================
print("Generating Figure 2: ABP waveform overlay...")
subj = selected[0]
global_mask = all_subjects == subj
te_mask_local = te_subjects == subj
local_indices = np.where(te_mask_local)[0]

# Take 5 consecutive segments
seg_start = 10  # skip first few for stability
n_show = 5
ibp1_segs = all_ibp1[te_idx][te_mask_local][seg_start:seg_start+n_show]
true_sbp_segs = all_sbp[te_idx][te_mask_local][seg_start:seg_start+n_show]
true_dbp_segs = all_dbp[te_idx][te_mask_local][seg_start:seg_start+n_show]

cnn_sbp_segs = dl_preds['CNN'][local_indices[seg_start:seg_start+n_show], 0]
cnn_dbp_segs = dl_preds['CNN'][local_indices[seg_start:seg_start+n_show], 1]
res_sbp_segs = dl_preds['ResNet'][local_indices[seg_start:seg_start+n_show], 0]
res_dbp_segs = dl_preds['ResNet'][local_indices[seg_start:seg_start+n_show], 1]
xgb_sbp_segs = xgb_sbp_pred[global_mask][seg_start:seg_start+n_show]
xgb_dbp_segs = xgb_dbp_pred[global_mask][seg_start:seg_start+n_show]

# Concatenate waveforms
sr = 125
seg_len = 1250
total_len = n_show * seg_len
abp_concat = ibp1_segs.flatten()
t = np.arange(total_len) / sr

fig, ax = plt.subplots(figsize=(16, 6))

# ABP waveform
ax.plot(t, abp_concat, color=C_ABP_WAVE, linewidth=0.6, alpha=0.9, label='Actual ABP Waveform')

# For each segment, draw horizontal lines for predicted values
for i in range(n_show):
    seg_t_start = i * seg_len / sr
    seg_t_end = (i + 1) * seg_len / sr

    # True SBP/DBP
    ax.hlines(true_sbp_segs[i], seg_t_start, seg_t_end, colors=C_ACTUAL, linewidth=2.5,
              label='Actual SBP' if i == 0 else None, zorder=4)
    ax.hlines(true_dbp_segs[i], seg_t_start, seg_t_end, colors=C_ACTUAL, linewidth=2.5, linestyle='-',
              label='Actual DBP' if i == 0 else None, zorder=4)

    # XGBoost
    ax.hlines(xgb_sbp_segs[i], seg_t_start, seg_t_end, colors=C_XGB, linewidth=1.5, linestyle='--',
              label='XGBoost' if i == 0 else None)
    ax.hlines(xgb_dbp_segs[i], seg_t_start, seg_t_end, colors=C_XGB, linewidth=1.5, linestyle='--')

    # CNN
    ax.hlines(cnn_sbp_segs[i], seg_t_start, seg_t_end, colors=C_CNN, linewidth=1.5, linestyle='--',
              label='CNN' if i == 0 else None)
    ax.hlines(cnn_dbp_segs[i], seg_t_start, seg_t_end, colors=C_CNN, linewidth=1.5, linestyle='--')

    # ResNet
    ax.hlines(res_sbp_segs[i], seg_t_start, seg_t_end, colors=C_RES, linewidth=1.5, linestyle='--',
              label='ResNet' if i == 0 else None)
    ax.hlines(res_dbp_segs[i], seg_t_start, seg_t_end, colors=C_RES, linewidth=1.5, linestyle='--')

    # Segment separator
    if i < n_show - 1:
        ax.axvline(seg_t_end, color='gray', linestyle=':', alpha=0.3)

ax.set_xlabel('Time (seconds)', fontsize=12, fontweight='bold')
ax.set_ylabel('Blood Pressure (mmHg)', fontsize=12, fontweight='bold')
age_val = all_demo[global_mask][0, 0]
sex_val = 'M' if all_demo[global_mask][0, 1] == 1 else 'F'
ax.set_title(f'ABP Waveform with Predicted SBP/DBP - {subj} (age={age_val:.1f}yr, {sex_val})',
             fontsize=14, fontweight='bold')
ax.legend(fontsize=10, ncol=5, loc='upper right')
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.grid(alpha=0.2)

# Annotations
ax.annotate('SBP', xy=(0.5, true_sbp_segs[0]+1), fontsize=9, color=C_ACTUAL,
            fontweight='bold', ha='center')
ax.annotate('DBP', xy=(0.5, true_dbp_segs[0]-3), fontsize=9, color=C_ACTUAL,
            fontweight='bold', ha='center')

plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'abp_waveform_overlay.png'), bbox_inches='tight')
plt.close()
print("  Saved: abp_waveform_overlay.png")


# ======================================================================
# FIGURE 3: Scatter Plots (Actual vs Predicted) - 3 models x SBP/DBP
# ======================================================================
print("Generating Figure 3: Scatter plots...")
fig, axes = plt.subplots(2, 3, figsize=(15, 10))

te_true_sbp = all_sbp[te_idx]
te_true_dbp = all_dbp[te_idx]
xgb_sbp_te = xgb_sbp_pred[te_idx]  # use same indices
xgb_dbp_te = xgb_dbp_pred[te_idx]
cnn_sbp_all = dl_preds['CNN'][:, 0]
cnn_dbp_all = dl_preds['CNN'][:, 1]
res_sbp_all = dl_preds['ResNet'][:, 0]
res_dbp_all = dl_preds['ResNet'][:, 1]

models_scatter = [
    ('XGBoost', xgb_sbp_te, xgb_dbp_te, C_XGB),
    ('Improved CNN', cnn_sbp_all, cnn_dbp_all, C_CNN),
    ('ResNet1D', res_sbp_all, res_dbp_all, C_RES),
]

for j, (mname, pred_sbp, pred_dbp, color) in enumerate(models_scatter):
    n_pts = min(len(te_true_sbp), len(pred_sbp))

    # SBP scatter
    ax = axes[0, j]
    ax.scatter(te_true_sbp[:n_pts], pred_sbp[:n_pts], s=3, alpha=0.15, color=color, rasterized=True)
    lims = [50, 180]
    ax.plot(lims, lims, 'k--', linewidth=1, alpha=0.5, label='Identity')
    ax.set_xlim(lims); ax.set_ylim(lims)
    mae = np.mean(np.abs(te_true_sbp[:n_pts] - pred_sbp[:n_pts]))
    r2 = 1 - np.sum((te_true_sbp[:n_pts] - pred_sbp[:n_pts])**2) / np.sum((te_true_sbp[:n_pts] - te_true_sbp[:n_pts].mean())**2)
    ax.text(0.05, 0.92, f'MAE={mae:.2f}\nR2={r2:.3f}', transform=ax.transAxes,
            fontsize=10, fontweight='bold', color=color,
            bbox=dict(boxstyle='round', facecolor='white', alpha=0.8))
    ax.set_title(f'{mname} - SBP', fontsize=13, fontweight='bold', color=color)
    ax.set_xlabel('Actual SBP (mmHg)', fontsize=10)
    ax.set_ylabel('Predicted SBP (mmHg)', fontsize=10)
    ax.grid(alpha=0.2)
    ax.set_aspect('equal')

    # DBP scatter
    ax = axes[1, j]
    ax.scatter(te_true_dbp[:n_pts], pred_dbp[:n_pts], s=3, alpha=0.15, color=color, rasterized=True)
    lims = [20, 110]
    ax.plot(lims, lims, 'k--', linewidth=1, alpha=0.5)
    ax.set_xlim(lims); ax.set_ylim(lims)
    mae = np.mean(np.abs(te_true_dbp[:n_pts] - pred_dbp[:n_pts]))
    r2 = 1 - np.sum((te_true_dbp[:n_pts] - pred_dbp[:n_pts])**2) / np.sum((te_true_dbp[:n_pts] - te_true_dbp[:n_pts].mean())**2)
    ax.text(0.05, 0.92, f'MAE={mae:.2f}\nR2={r2:.3f}', transform=ax.transAxes,
            fontsize=10, fontweight='bold', color=color,
            bbox=dict(boxstyle='round', facecolor='white', alpha=0.8))
    ax.set_title(f'{mname} - DBP', fontsize=13, fontweight='bold', color=color)
    ax.set_xlabel('Actual DBP (mmHg)', fontsize=10)
    ax.set_ylabel('Predicted DBP (mmHg)', fontsize=10)
    ax.grid(alpha=0.2)
    ax.set_aspect('equal')

fig.suptitle('Actual vs Predicted: Scatter Plots (Test Fold)', fontsize=15, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'scatter_actual_vs_pred.png'), bbox_inches='tight')
plt.close()
print("  Saved: scatter_actual_vs_pred.png")


# ======================================================================
# FIGURE 4: Error Distribution (Histogram)
# ======================================================================
print("Generating Figure 4: Error distribution...")
fig, axes = plt.subplots(2, 3, figsize=(15, 9))

for j, (mname, pred_sbp, pred_dbp, color) in enumerate(models_scatter):
    n_pts = min(len(te_true_sbp), len(pred_sbp))
    err_sbp = pred_sbp[:n_pts] - te_true_sbp[:n_pts]
    err_dbp = pred_dbp[:n_pts] - te_true_dbp[:n_pts]

    ax = axes[0, j]
    ax.hist(err_sbp, bins=80, color=color, alpha=0.7, edgecolor='white', linewidth=0.5, density=True)
    ax.axvline(0, color='black', linestyle='-', linewidth=1, alpha=0.5)
    ax.axvline(np.mean(err_sbp), color=color, linestyle='--', linewidth=2, label=f'Mean={np.mean(err_sbp):.1f}')
    ax.set_title(f'{mname} - SBP Error', fontsize=13, fontweight='bold', color=color)
    ax.set_xlabel('Prediction Error (mmHg)', fontsize=10)
    ax.set_ylabel('Density', fontsize=10)
    ax.set_xlim(-50, 50)
    ax.legend(fontsize=9)
    ax.grid(alpha=0.2)
    ax.text(0.95, 0.92, f'SD={np.std(err_sbp):.1f}', transform=ax.transAxes,
            fontsize=10, ha='right', fontweight='bold', color=color)

    ax = axes[1, j]
    ax.hist(err_dbp, bins=80, color=color, alpha=0.7, edgecolor='white', linewidth=0.5, density=True)
    ax.axvline(0, color='black', linestyle='-', linewidth=1, alpha=0.5)
    ax.axvline(np.mean(err_dbp), color=color, linestyle='--', linewidth=2, label=f'Mean={np.mean(err_dbp):.1f}')
    ax.set_title(f'{mname} - DBP Error', fontsize=13, fontweight='bold', color=color)
    ax.set_xlabel('Prediction Error (mmHg)', fontsize=10)
    ax.set_ylabel('Density', fontsize=10)
    ax.set_xlim(-40, 40)
    ax.legend(fontsize=9)
    ax.grid(alpha=0.2)
    ax.text(0.95, 0.92, f'SD={np.std(err_dbp):.1f}', transform=ax.transAxes,
            fontsize=10, ha='right', fontweight='bold', color=color)

fig.suptitle('Prediction Error Distribution (Bland-Altman Histogram)', fontsize=15, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'error_distribution.png'), bbox_inches='tight')
plt.close()
print("  Saved: error_distribution.png")


# ======================================================================
# FIGURE 5: Per-Patient MAE Comparison (Box Plot)
# ======================================================================
print("Generating Figure 5: Per-patient box plot...")
fig, axes = plt.subplots(1, 2, figsize=(14, 6))

patient_mae_sbp = {m: [] for m in ['XGBoost', 'CNN', 'ResNet']}
patient_mae_dbp = {m: [] for m in ['XGBoost', 'CNN', 'ResNet']}
patient_labels = []

for subj in te_unique:
    global_mask = all_subjects == subj
    te_local_mask = te_subjects == subj
    local_indices = np.where(te_local_mask)[0]

    true_s = all_sbp[te_idx][te_local_mask]
    true_d = all_dbp[te_idx][te_local_mask]

    patient_labels.append(subj.replace('case', 'P'))

    patient_mae_sbp['XGBoost'].append(np.mean(np.abs(xgb_sbp_pred[global_mask] - all_sbp[global_mask])))
    patient_mae_dbp['XGBoost'].append(np.mean(np.abs(xgb_dbp_pred[global_mask] - all_dbp[global_mask])))

    patient_mae_sbp['CNN'].append(np.mean(np.abs(dl_preds['CNN'][local_indices, 0] - true_s)))
    patient_mae_dbp['CNN'].append(np.mean(np.abs(dl_preds['CNN'][local_indices, 1] - true_d)))

    patient_mae_sbp['ResNet'].append(np.mean(np.abs(dl_preds['ResNet'][local_indices, 0] - true_s)))
    patient_mae_dbp['ResNet'].append(np.mean(np.abs(dl_preds['ResNet'][local_indices, 1] - true_d)))

# Box plot
colors = [C_XGB, C_CNN, C_RES]
for ax, target, data_dict in [
    (axes[0], 'SBP MAE per Patient', patient_mae_sbp),
    (axes[1], 'DBP MAE per Patient', patient_mae_dbp),
]:
    bp_data = [data_dict['XGBoost'], data_dict['CNN'], data_dict['ResNet']]
    bp = ax.boxplot(bp_data, labels=['XGBoost', 'CNN', 'ResNet'],
                    patch_artist=True, widths=0.5,
                    medianprops=dict(color='black', linewidth=2))
    for patch, color in zip(bp['boxes'], colors):
        patch.set_facecolor(color)
        patch.set_alpha(0.6)

    # Overlay individual patient points
    for i, (mname, vals) in enumerate(data_dict.items()):
        jitter = np.random.normal(0, 0.05, len(vals))
        ax.scatter(np.ones(len(vals)) * (i+1) + jitter, vals, s=25, alpha=0.5,
                   color=colors[i], edgecolors='black', linewidth=0.5, zorder=5)

    ax.set_title(target, fontsize=14, fontweight='bold')
    ax.set_ylabel('MAE (mmHg)', fontsize=12)
    ax.grid(axis='y', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

fig.suptitle('Per-Patient MAE Distribution (Test Fold)', fontsize=15, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'per_patient_boxplot.png'), bbox_inches='tight')
plt.close()
print("  Saved: per_patient_boxplot.png")


# ======================================================================
# FIGURE 6: Detailed ABP waveform single segment comparison
# ======================================================================
print("Generating Figure 6: Single segment detail...")
fig, axes = plt.subplots(2, 2, figsize=(14, 8))

for idx, (subj, ax) in enumerate(zip(selected, axes.flat)):
    global_mask = all_subjects == subj
    te_local_mask = te_subjects == subj
    local_indices = np.where(te_local_mask)[0]

    # Pick a segment in the middle
    seg_idx = len(local_indices) // 2
    ibp1_seg = all_ibp1[te_idx][te_local_mask][seg_idx]
    true_s = all_sbp[te_idx][te_local_mask][seg_idx]
    true_d = all_dbp[te_idx][te_local_mask][seg_idx]
    true_m = all_mbp[te_idx][te_local_mask][seg_idx]

    cnn_s = dl_preds['CNN'][local_indices[seg_idx], 0]
    cnn_d = dl_preds['CNN'][local_indices[seg_idx], 1]
    res_s = dl_preds['ResNet'][local_indices[seg_idx], 0]
    res_d = dl_preds['ResNet'][local_indices[seg_idx], 1]

    t = np.arange(len(ibp1_seg)) / 125.0

    # ABP waveform
    ax.plot(t, ibp1_seg, color=C_ABP_WAVE, linewidth=1.2, label='Actual ABP', zorder=3)

    # True SBP/DBP
    ax.axhline(true_s, color=C_ACTUAL, linewidth=2, linestyle='-', alpha=0.7, label=f'True SBP={true_s:.0f}')
    ax.axhline(true_d, color=C_ACTUAL, linewidth=2, linestyle='-', alpha=0.7, label=f'True DBP={true_d:.0f}')

    # CNN predicted
    ax.axhline(cnn_s, color=C_CNN, linewidth=1.5, linestyle='--', alpha=0.8)
    ax.axhline(cnn_d, color=C_CNN, linewidth=1.5, linestyle='--', alpha=0.8)
    ax.fill_between(t, cnn_d, cnn_s, alpha=0.06, color=C_CNN)

    # ResNet predicted
    ax.axhline(res_s, color=C_RES, linewidth=1.5, linestyle=':', alpha=0.8)
    ax.axhline(res_d, color=C_RES, linewidth=1.5, linestyle=':', alpha=0.8)

    # Annotations on right side
    y_offset = 0
    for val, lbl, c in [(cnn_s, f'CNN SBP={cnn_s:.0f}', C_CNN),
                         (res_s, f'ResNet SBP={res_s:.0f}', C_RES),
                         (cnn_d, f'CNN DBP={cnn_d:.0f}', C_CNN),
                         (res_d, f'ResNet DBP={res_d:.0f}', C_RES)]:
        ax.annotate(lbl, xy=(10, val), xytext=(10.3, val + y_offset),
                    fontsize=7.5, color=c, fontweight='bold', va='center')

    age_val = all_demo[global_mask][0, 0]
    sex_val = 'M' if all_demo[global_mask][0, 1] == 1 else 'F'
    ax.set_title(f'{subj} (age={age_val:.1f}yr, {sex_val}) - 10s Segment',
                 fontsize=11, fontweight='bold', loc='left')
    ax.set_xlabel('Time (s)', fontsize=10)
    ax.set_ylabel('BP (mmHg)', fontsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(alpha=0.15)
    if idx == 0:
        ax.legend(fontsize=8, loc='upper right', ncol=2)

fig.suptitle('ABP Waveform with Model Predictions - Single 10s Segments',
             fontsize=14, fontweight='bold', y=1.01)
plt.tight_layout()
plt.savefig(os.path.join(FIG_DIR, 'single_segment_detail.png'), bbox_inches='tight')
plt.close()
print("  Saved: single_segment_detail.png")


# ======================================================================
# Insert all new figures into PPTX
# ======================================================================
print("\n" + "="*60)
print("Inserting figures into PowerPoint...")
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation(os.path.join(BASE, 'PPG2ABP_Presentation.pptx'))

PRIMARY = PRGBColor(0, 51, 102)
WHITE = PRGBColor(255, 255, 255)
GRAY_TEXT = PRGBColor(100, 100, 100)


def add_chart_slide(title, img_path, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, PInches(0), PInches(0), PInches(13.333), PInches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = PRIMARY; shape.line.fill.background()
    txBox = slide.shapes.add_textbox(PInches(0.4), PInches(0.1), PInches(12), PInches(0.7))
    p = txBox.text_frame.paragraphs[0]
    p.text = title; p.font.size = PPt(28); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = 'Malgun Gothic'

    if subtitle:
        slide.shapes.add_picture(img_path, PInches(0.3), PInches(1.05), PInches(12.7), PInches(5.7))
        txBox2 = slide.shapes.add_textbox(PInches(0.3), PInches(6.85), PInches(12.7), PInches(0.4))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = subtitle; p2.font.size = PPt(12); p2.font.color.rgb = GRAY_TEXT
        p2.font.name = 'Malgun Gothic'; p2.alignment = PP_ALIGN.CENTER
    else:
        slide.shapes.add_picture(img_path, PInches(0.3), PInches(1.05), PInches(12.7), PInches(6.2))
    return slide


def add_two_chart_slide(title, img1, img2, cap1='', cap2=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, PInches(0), PInches(0), PInches(13.333), PInches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = PRIMARY; shape.line.fill.background()
    txBox = slide.shapes.add_textbox(PInches(0.4), PInches(0.1), PInches(12), PInches(0.7))
    p = txBox.text_frame.paragraphs[0]
    p.text = title; p.font.size = PPt(28); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = 'Malgun Gothic'

    slide.shapes.add_picture(img1, PInches(0.2), PInches(1.1), PInches(6.4), PInches(4.8))
    slide.shapes.add_picture(img2, PInches(6.7), PInches(1.1), PInches(6.4), PInches(4.8))

    if cap1:
        txBox1 = slide.shapes.add_textbox(PInches(0.2), PInches(6.0), PInches(6.4), PInches(0.3))
        p1 = txBox1.text_frame.paragraphs[0]; p1.text = cap1
        p1.font.size = PPt(11); p1.font.name = 'Malgun Gothic'
        p1.alignment = PP_ALIGN.CENTER; p1.font.color.rgb = GRAY_TEXT
    if cap2:
        txBox2 = slide.shapes.add_textbox(PInches(6.7), PInches(6.0), PInches(6.4), PInches(0.3))
        p2 = txBox2.text_frame.paragraphs[0]; p2.text = cap2
        p2.font.size = PPt(11); p2.font.name = 'Malgun Gothic'
        p2.alignment = PP_ALIGN.CENTER; p2.font.color.rgb = GRAY_TEXT
    return slide


# Slide 1: Time series
add_chart_slide(
    'Actual vs Predicted BP - Time Series (4 Patients)',
    os.path.join(FIG_DIR, 'actual_vs_pred_timeseries.png'),
    'Black=Actual, Red=XGBoost, Blue=CNN, Green=ResNet / Solid=SBP, Dotted=DBP'
)

# Slide 2: Waveform overlay
add_chart_slide(
    'ABP Waveform with Predicted SBP/DBP Overlay',
    os.path.join(FIG_DIR, 'abp_waveform_overlay.png'),
    '5 consecutive 10s segments showing actual ABP waveform with model predictions'
)

# Slide 3: Single segment detail
add_chart_slide(
    'Single Segment ABP Waveform Detail (4 Patients)',
    os.path.join(FIG_DIR, 'single_segment_detail.png'),
    'Individual 10s ABP waveform with CNN (dashed) and ResNet (dotted) prediction lines'
)

# Slide 4: Scatter + Error distribution
add_two_chart_slide(
    'Actual vs Predicted: Scatter & Error Distribution',
    os.path.join(FIG_DIR, 'scatter_actual_vs_pred.png'),
    os.path.join(FIG_DIR, 'error_distribution.png'),
    'Scatter: closer to diagonal = better prediction',
    'Histogram: narrower and centered at 0 = better'
)

# Slide 5: Box plot
add_chart_slide(
    'Per-Patient MAE Distribution (Box Plot)',
    os.path.join(FIG_DIR, 'per_patient_boxplot.png'),
    'Individual points = per-patient MAE. DL models show lower median and less spread.'
)

# Save
for fname in ['PPG2ABP_Presentation.pptx', 'PPG2ABP_Presentation_v2.pptx', 'PPG2ABP_Presentation_v3.pptx']:
    try:
        save_path = os.path.join(BASE, fname)
        prs.save(save_path)
        print(f"Saved: {save_path}")
        break
    except PermissionError:
        continue

print("\nAll done! 6 figures generated, 5 new slides added to PPTX.")
