"""
Update Word and PPT reports with Step 4 GPU model results.
"""
import numpy as np
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = 'C:/Users/jaege/Desktop/Study/PPG2ABP'

# ========================================================================
# PART 1: UPDATE WORD DOCUMENT
# ========================================================================
print("Updating Word report...")
doc = Document(os.path.join(BASE, 'PPG2ABP_Report.docx'))

# Helper
def add_heading_styled(text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0, 51, 102)
    return h

def add_table_with_style(headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = 'Light Grid Accent 1'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for j, h in enumerate(headers):
        cell = table.rows[0].cells[j]
        cell.text = h
        for p in cell.paragraphs:
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in p.runs:
                run.bold = True
                run.font.size = Pt(10)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.rows[i+1].cells[j]
            cell.text = str(val)
            for p in cell.paragraphs:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in p.runs:
                    run.font.size = Pt(10)
    return table

# ── Add new chapter before Conclusion ──
# Find the paragraph index of "6. 결론"
conclusion_idx = None
for i, para in enumerate(doc.paragraphs):
    if '6. 결론' in para.text or '6. 결론' in para.text:
        conclusion_idx = i
        break

# We'll add content at the end (before references), since inserting at specific
# positions in python-docx is complex. Instead, add as new sections after existing content.
# First, let's add a page break before the new content.

# Renumber: existing 6.결론 -> 7, 7.참고문헌 -> 8
# Add Step 4 as new section 5 (between 4.결과 and old 5.고찰)

# Since modifying existing paragraphs is complex, we append new chapters at the end
doc.add_page_break()

# ════════════════════════════════════════════════════════════════
# NEW CHAPTER: Step 4 GPU Models
# ════════════════════════════════════════════════════════════════
add_heading_styled('부록 A. GPU 기반 확장 모델 (Step 4)', level=1)

doc.add_paragraph(
    '초기 분석 이후 PyTorch CUDA 환경을 구축하여 (RTX 4070 Laptop, 8GB VRAM, '
    'PyTorch 2.6.0+cu124) 전체 데이터(81,812 샘플)에 대해 더 강력한 딥러닝 모델을 '
    '학습하였다. 5-Fold Group K-Fold 교차 검증을 적용하였으며, 환자 단위 분리를 유지하였다.'
)

add_heading_styled('A.1 Improved 1D-CNN', level=2)
doc.add_paragraph(
    '기존 경량 CNN을 확장하여 5개의 합성곱 블록으로 구성된 깊은 네트워크를 설계하였다. '
    '각 블록은 2개의 Conv1d-BatchNorm-ReLU 층과 MaxPool, Dropout으로 구성된다. '
    '최종적으로 512차원의 Global Average Pooling 특징에 나이/성별을 concat하여 '
    '3개의 출력(SBP, DBP, MBP)을 예측한다.'
)

doc.add_paragraph('주요 학습 설정:')
settings = [
    'Optimizer: AdamW (weight_decay=1e-4)',
    'Scheduler: CosineAnnealingLR',
    'Loss: Smooth L1 Loss (Huber)',
    'Batch size: 512, Max epochs: 40, Early stopping patience: 8',
    'Gradient clipping: max_norm=1.0',
]
for s in settings:
    doc.add_paragraph(s, style='List Bullet')

add_heading_styled('A.2 ResNet1D', level=2)
doc.add_paragraph(
    'ResNet 아키텍처를 1D 시계열에 적용하였다. Stem layer (Conv1d, stride=2, MaxPool) 이후 '
    '4개의 Residual Stage (각 2개의 ResBlock)로 구성되며, 채널 수는 64-128-256-512로 '
    '점진적으로 증가한다. 각 ResBlock은 skip connection을 포함하여 gradient flow를 개선한다. '
    'Global Average Pooling 후 나이/성별을 concat하여 예측한다.'
)

add_heading_styled('A.3 U-Net 1D (파형 재구성)', level=2)
doc.add_paragraph(
    'PPG 파형으로부터 ABP 파형 전체를 재구성하는 1D U-Net 모델을 설계하였다. '
    'Encoder (4단계 다운샘플링), Bottleneck, Decoder (4단계 업샘플링 + skip connection)로 '
    '구성되며, Bottleneck에서 FiLM (Feature-wise Linear Modulation) 기법으로 나이/성별 '
    '정보를 주입한다. 손실 함수는 MSE와 waveform gradient loss의 가중 합을 사용하여 '
    '파형의 형태(shape)를 보존하도록 하였다.'
)
doc.add_paragraph(
    'FiLM conditioning: 인구통계 정보(나이, 성별)로부터 scale(gamma)과 shift(beta) 파라미터를 '
    '학습하여 bottleneck feature map에 affine transformation을 적용한다.'
)
doc.add_paragraph(
    'U-Net은 5개 fold 중 2개만 완료되었으며 (fold당 약 45분 소요), '
    '완료된 2개 fold의 평균 파형 상관계수는 0.68이었다.'
)

doc.add_page_break()

add_heading_styled('A.4 GPU 모델 결과', level=2)

add_heading_styled('표 A1. Improved 1D-CNN 결과 (5-Fold Group CV, 전체 데이터)', level=3)
add_table_with_style(
    ['Target', 'MAE\n(mmHg)', 'RMSE\n(mmHg)', 'R2', 'Bias+/-SD\n(mmHg)',
     '<=5\nmmHg', '<=10\nmmHg', '<=15\nmmHg', 'BHS'],
    [
        ['SBP', '12.64', '16.21', '0.1731', '-0.9+/-16.2', '25.3%', '48.1%', '66.9%', 'D'],
        ['DBP', '7.20', '9.83', '0.2282', '-0.7+/-9.8', '44.8%', '76.7%', '90.0%', 'C'],
        ['MBP', '8.08', '11.14', '0.2717', '-0.9+/-11.1', '42.5%', '70.8%', '86.1%', 'C'],
    ])
doc.add_paragraph()

add_heading_styled('표 A2. Improved 1D-CNN Fold별 성능', level=3)
add_table_with_style(
    ['Fold', 'Test\nSubjects', 'Test\nSamples', 'SBP MAE', 'DBP MAE', 'MBP MAE', 'Epochs', 'Time(s)'],
    [
        ['1', '9', '16,438', '12.10', '7.72', '8.89', '18', '411'],
        ['2', '10', '16,255', '12.36', '6.65', '7.65', '17', '388'],
        ['3', '9', '16,419', '11.31', '6.66', '7.36', '25', '570'],
        ['4', '9', '16,280', '13.66', '7.89', '8.57', '17', '380'],
        ['5', '9', '16,420', '13.77', '7.08', '7.93', '10', '231'],
    ])
doc.add_paragraph()

add_heading_styled('표 A3. ResNet1D 결과 (5-Fold Group CV, 전체 데이터)', level=3)
add_table_with_style(
    ['Target', 'MAE\n(mmHg)', 'RMSE\n(mmHg)', 'R2', 'Bias+/-SD\n(mmHg)',
     '<=5\nmmHg', '<=10\nmmHg', '<=15\nmmHg', 'BHS'],
    [
        ['SBP', '12.25', '15.91', '0.2031', '-2.2+/-15.8', '27.5%', '51.9%', '67.9%', 'D'],
        ['DBP', '7.83', '10.18', '0.1724', '-0.5+/-10.2', '39.0%', '71.7%', '89.2%', 'D'],
        ['MBP', '8.22', '11.11', '0.2763', '-0.8+/-11.1', '40.2%', '70.5%', '85.0%', 'C'],
    ])
doc.add_paragraph()

add_heading_styled('표 A4. ResNet1D Fold별 성능', level=3)
add_table_with_style(
    ['Fold', 'Test\nSubjects', 'Test\nSamples', 'SBP MAE', 'DBP MAE', 'MBP MAE', 'Epochs', 'Time(s)'],
    [
        ['1', '9', '16,438', '14.02', '9.09', '9.31', '10', '316'],
        ['2', '10', '16,255', '10.57', '7.49', '7.78', '22', '677'],
        ['3', '9', '16,419', '11.75', '7.51', '8.31', '13', '410'],
        ['4', '9', '16,280', '12.01', '7.28', '7.84', '23', '701'],
        ['5', '9', '16,420', '12.91', '7.76', '7.87', '13', '409'],
    ])
doc.add_paragraph()

add_heading_styled('표 A5. U-Net 1D 파형 재구성 결과 (부분 완료)', level=3)
add_table_with_style(
    ['Fold', 'Test Samples', 'Waveform RMSE', 'Waveform Correlation', 'Epochs', 'Time(s)'],
    [
        ['1', '16,438', '13.40', '0.587', '40', '2,669'],
        ['2', '16,255', '9.66', '0.772', '40', '32,469'],
        ['Mean', '-', '11.53', '0.680', '-', '-'],
    ])
doc.add_paragraph()

doc.add_page_break()

add_heading_styled('A.5 전체 모델 비교', level=2)
doc.add_paragraph(
    '전체 실험에서 사용된 6가지 모델의 성능을 비교한 결과를 아래 표에 정리하였다.'
)

add_heading_styled('표 A6. 전체 모델 성능 비교 종합', level=3)
add_table_with_style(
    ['Model', 'CV Method', 'Data\nSize', 'SBP\nMAE', 'DBP\nMAE', 'MBP\nMAE',
     'SBP\nR2', 'DBP\nBHS', 'MBP\nBHS'],
    [
        ['XGBoost', 'LOSO (46)', '81,812', '14.98', '8.66', '10.41', '-0.19', 'D', 'D'],
        ['LightGBM', 'LOSO (46)', '81,812', '16.09', '8.42', '10.15', '-0.34', 'D', 'D'],
        ['Light CNN\n(CPU, sub)', '5-Fold', '22,638', '12.95', '7.47', '8.49', '0.16', 'C', 'D'],
        ['Improved CNN\n(GPU, full)', '5-Fold', '81,812', '12.64', '7.20', '8.08', '0.17', 'C', 'C'],
        ['ResNet1D\n(GPU, full)', '5-Fold', '81,812', '12.25', '7.83', '8.22', '0.20', 'D', 'C'],
        ['U-Net 1D\n(GPU, partial)', '5-Fold', '81,812', '-', '-', '-', '-', 'Corr\n=0.68', '-'],
    ])
doc.add_paragraph()

doc.add_paragraph(
    '주요 발견: (1) Improved CNN이 전체적으로 가장 균형 잡힌 성능을 보였으며, DBP와 MBP '
    '모두 BHS Grade C를 달성하였다. (2) ResNet1D는 SBP 예측에서 최고 성능(MAE 12.25, '
    'R2 0.20)을 보였으나 DBP에서는 CNN보다 낮은 성능을 보였다. (3) 전체 데이터(81,812)를 '
    '사용한 경우 서브샘플(22,638) 대비 CNN의 DBP MAE가 7.47에서 7.20으로 개선되었다. '
    '(4) U-Net은 PPG 파형에서 ABP 파형을 일정 수준(Correlation 0.68) 재구성할 수 있음을 보였다.'
)

add_heading_styled('A.6 모델별 장단점 분석', level=2)

doc.add_paragraph('Improved 1D-CNN:', style='Heading 4')
pros_cons = [
    '장점: DBP/MBP 모두 BHS Grade C, 가장 낮은 DBP MAE (7.20), 학습 시간 33분',
    '단점: SBP 예측력 ResNet1D 대비 약간 낮음',
    '추천: 전반적 혈압 예측이 필요한 경우 최적의 선택',
]
for p in pros_cons:
    doc.add_paragraph(p, style='List Bullet')

doc.add_paragraph('ResNet1D:', style='Heading 4')
pros_cons = [
    '장점: SBP 최저 MAE (12.25), 최고 SBP R2 (0.20), residual connection으로 안정적 학습',
    '단점: DBP 성능이 CNN보다 낮음 (MAE 7.83 vs 7.20), 학습 시간 42분',
    '추천: SBP 예측이 중요한 경우 적합',
]
for p in pros_cons:
    doc.add_paragraph(p, style='List Bullet')

doc.add_paragraph('U-Net 1D:', style='Heading 4')
pros_cons = [
    '장점: ABP 파형 전체를 재구성, 풍부한 혈역학적 정보 제공, FiLM conditioning',
    '단점: 학습 시간이 매우 김 (fold당 45분~9시간), 완전한 평가 미완료',
    '추천: 파형 형태가 중요한 연구 목적에 적합, GPU 자원이 충분한 환경 필요',
]
for p in pros_cons:
    doc.add_paragraph(p, style='List Bullet')

add_heading_styled('A.7 GPU 환경 정보', level=2)
add_table_with_style(
    ['항목', '사양'],
    [
        ['GPU', 'NVIDIA GeForce RTX 4070 Laptop'],
        ['VRAM', '8 GB'],
        ['PyTorch', '2.6.0+cu124'],
        ['CUDA', '12.4'],
        ['Driver', '560.94'],
        ['CNN 총 학습 시간', '1,981초 (33분)'],
        ['ResNet1D 총 학습 시간', '2,513초 (42분)'],
        ['U-Net fold당 학습 시간', '2,669 ~ 32,469초'],
    ])

# Save
doc.save(os.path.join(BASE, 'PPG2ABP_Report.docx'))
print("Word report updated.")


# ========================================================================
# PART 2: UPDATE POWERPOINT
# ========================================================================
print("Updating PowerPoint...")
prs = Presentation(os.path.join(BASE, 'PPG2ABP_Presentation.pptx'))

# Color scheme
PRIMARY = PRGBColor(0, 51, 102)
SECONDARY = PRGBColor(0, 102, 153)
ACCENT = PRGBColor(220, 80, 60)
LIGHT_BG = PRGBColor(240, 245, 250)
WHITE = PRGBColor(255, 255, 255)
BLACK = PRGBColor(40, 40, 40)
GRAY = PRGBColor(120, 120, 120)
GREEN = PRGBColor(40, 160, 80)
TABLE_HEADER = PRGBColor(0, 70, 130)
TABLE_ALT = PRGBColor(230, 240, 250)

def add_slide(layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PPt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    PInches(left), PInches(top), PInches(width), PInches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_title_bar(slide, title_text):
    add_rect(slide, 0, 0, 13.333, 1.1, PRIMARY)
    add_textbox(slide, 0.5, 0.15, 12, 0.8, title_text, 32, True, WHITE, PP_ALIGN.LEFT)

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=BLACK):
    txBox = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = PPt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Malgun Gothic'
        p.space_after = PPt(4)
    return txBox

def add_table_to_slide(slide, left, top, width, height, headers, rows):
    n_rows, n_cols = len(rows) + 1, len(headers)
    ts = slide.shapes.add_table(n_rows, n_cols, PInches(left), PInches(top),
                                 PInches(width), PInches(height))
    table = ts.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_HEADER
        for p in cell.text_frame.paragraphs:
            p.font.size = PPt(11); p.font.bold = True
            p.font.color.rgb = WHITE; p.font.name = 'Malgun Gothic'
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j); cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_ALT
            for p in cell.text_frame.paragraphs:
                p.font.size = PPt(10); p.font.color.rgb = BLACK
                p.font.name = 'Malgun Gothic'; p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return ts

# ── Insert new slides before "Thank You" (last slide) ──
# We'll add slides at the end, then the Thank You slide is already there.
# Actually, let's just add new slides. The Thank You is the last slide.

# SLIDE: GPU Models Architecture
slide = add_slide()
add_title_bar(slide, 'Step 4. GPU 확장 모델 - 아키텍처')

# Improved CNN
add_rect(slide, 0.3, 1.3, 4.1, 2.8, PRGBColor(230, 245, 255))
add_textbox(slide, 0.5, 1.35, 3.8, 0.4, 'Improved 1D-CNN', 18, True, PRIMARY)
add_bullet_text(slide, 0.5, 1.85, 3.8, 2.2, [
    '5 Conv blocks (double Conv each)',
    'Channels: 32-64-128-256-512',
    'Kernel: 15-9-5-3-3 (decreasing)',
    'GAP + concat [age, sex]',
    'FC: 514->256->64->3',
    'Total: ~1.5M params',
], 13, BLACK)

# ResNet1D
add_rect(slide, 4.6, 1.3, 4.1, 2.8, PRGBColor(255, 245, 230))
add_textbox(slide, 4.8, 1.35, 3.8, 0.4, 'ResNet1D', 18, True, PRGBColor(180, 100, 0))
add_bullet_text(slide, 4.8, 1.85, 3.8, 2.2, [
    'Stem: Conv(k=15, stride=2) + Pool',
    '4 Residual Stages (2 blocks each)',
    'Channels: 64-128-256-512',
    'Skip connections (identity/1x1)',
    'GAP + concat [age, sex]',
    'Total: ~2.8M params',
], 13, BLACK)

# U-Net
add_rect(slide, 8.9, 1.3, 4.1, 2.8, PRGBColor(230, 255, 235))
add_textbox(slide, 9.1, 1.35, 3.8, 0.4, 'U-Net 1D', 18, True, GREEN)
add_bullet_text(slide, 9.1, 1.85, 3.8, 2.2, [
    'Encoder: 4 stages (32-64-128-256)',
    'Bottleneck: 512ch + FiLM',
    'Decoder: 4 stages + skip conn.',
    'FiLM: age/sex -> gamma, beta',
    'Loss: MSE + gradient loss',
    'Output: ABP waveform (1250 pts)',
], 13, BLACK)

# Training config
add_rect(slide, 0.3, 4.4, 12.7, 2.8, LIGHT_BG)
add_textbox(slide, 0.5, 4.5, 6, 0.4, 'GPU Training Config', 18, True, PRIMARY)
add_table_to_slide(slide, 0.5, 5.0, 6, 2,
    ['Setting', 'Value'],
    [
        ['GPU', 'RTX 4070 Laptop (8GB)'],
        ['PyTorch', '2.6.0 + CUDA 12.4'],
        ['Optimizer', 'AdamW (wd=1e-4)'],
        ['Scheduler', 'CosineAnnealingLR'],
        ['Batch / Epochs', '512 / 40 (patience=8)'],
        ['CV', '5-Fold Group (patient-level)'],
    ])

add_textbox(slide, 7, 5.0, 5.5, 0.4, 'Data', 18, True, PRIMARY)
add_bullet_text(slide, 7, 5.5, 5.5, 1.5, [
    'Full dataset: 81,812 segments',
    '46 subjects, 10s @ 125Hz each',
    'No subsampling (vs Step 3: 22,638)',
    'PPG z-score normalized per fold',
], 14, BLACK)


# SLIDE: CNN + ResNet Results
slide = add_slide()
add_title_bar(slide, 'Step 4. 결과 - Improved CNN & ResNet1D')

add_textbox(slide, 0.3, 1.3, 6.2, 0.4, 'Improved 1D-CNN (33min, GPU)', 18, True, PRIMARY)
add_table_to_slide(slide, 0.3, 1.8, 6.2, 2.2,
    ['Target', 'MAE', 'RMSE', 'R2', 'Bias+/-SD', '<=5', '<=10', '<=15', 'BHS'],
    [
        ['SBP', '12.64', '16.21', '0.17', '-0.9+/-16.2', '25.3%', '48.1%', '66.9%', 'D'],
        ['DBP', '7.20', '9.83', '0.23', '-0.7+/-9.8', '44.8%', '76.7%', '90.0%', 'C'],
        ['MBP', '8.08', '11.14', '0.27', '-0.9+/-11.1', '42.5%', '70.8%', '86.1%', 'C'],
    ])

add_textbox(slide, 6.8, 1.3, 6.2, 0.4, 'ResNet1D (42min, GPU)', 18, True, PRGBColor(180, 100, 0))
add_table_to_slide(slide, 6.8, 1.8, 6.2, 2.2,
    ['Target', 'MAE', 'RMSE', 'R2', 'Bias+/-SD', '<=5', '<=10', '<=15', 'BHS'],
    [
        ['SBP', '12.25', '15.91', '0.20', '-2.2+/-15.8', '27.5%', '51.9%', '67.9%', 'D'],
        ['DBP', '7.83', '10.18', '0.17', '-0.5+/-10.2', '39.0%', '71.7%', '89.2%', 'D'],
        ['MBP', '8.22', '11.11', '0.28', '-0.8+/-11.1', '40.2%', '70.5%', '85.0%', 'C'],
    ])

# Comparison boxes
add_rect(slide, 0.3, 4.3, 6.2, 2.8, PRGBColor(230, 245, 255))
add_textbox(slide, 0.5, 4.4, 5.8, 0.4, 'Improved CNN - Best For', 16, True, PRIMARY)
add_bullet_text(slide, 0.5, 4.9, 5.8, 2, [
    'DBP MAE 7.20 (All models best)',
    'DBP BHS Grade C + MBP BHS Grade C',
    'Best balanced across all targets',
    'Recommended as primary model',
], 14, BLACK)

add_rect(slide, 6.8, 4.3, 6.2, 2.8, PRGBColor(255, 245, 230))
add_textbox(slide, 7.0, 4.4, 5.8, 0.4, 'ResNet1D - Best For', 16, True, PRGBColor(180, 100, 0))
add_bullet_text(slide, 7.0, 4.9, 5.8, 2, [
    'SBP MAE 12.25 (All models best)',
    'SBP R2 0.20 (All models best)',
    'Skip connections help SBP learning',
    'Better when SBP accuracy matters',
], 14, BLACK)


# SLIDE: U-Net + Waveform
slide = add_slide()
add_title_bar(slide, 'Step 4. 결과 - U-Net 1D 파형 재구성')

add_textbox(slide, 0.3, 1.3, 6, 0.4, 'U-Net 1D (partial: 2/5 folds)', 18, True, GREEN)
add_table_to_slide(slide, 0.3, 1.8, 6, 1.8,
    ['Fold', 'Samples', 'WF RMSE', 'Correlation', 'Epochs', 'Time'],
    [
        ['1', '16,438', '13.40', '0.587', '40', '45min'],
        ['2', '16,255', '9.66', '0.772', '40', '~9hr'],
        ['Mean', '-', '11.53', '0.680', '-', '-'],
    ])

add_rect(slide, 0.3, 3.9, 6, 3.2, LIGHT_BG)
add_textbox(slide, 0.5, 4.0, 5.6, 0.4, 'U-Net Key Features', 16, True, GREEN)
add_bullet_text(slide, 0.5, 4.5, 5.6, 2.5, [
    'PPG waveform -> ABP waveform (1250 pts)',
    'FiLM conditioning (age/sex modulation)',
    'Combined loss: MSE + gradient loss',
    '  -> gradient loss preserves waveform shape',
    'Correlation 0.68 (avg of 2 folds)',
    '  -> Fold 2: 0.772 (promising!)',
    'Limitation: Very long training time',
    '  -> Needs more GPU resources or optimization',
], 13, BLACK)

add_rect(slide, 6.8, 1.3, 6.2, 5.8, PRGBColor(255, 250, 240))
add_textbox(slide, 7, 1.4, 5.8, 0.4, 'Waveform Reconstruction Concept', 16, True, PRGBColor(180, 100, 0))

add_textbox(slide, 7, 2.0, 5.8, 0.8,
    'Input: PPG (1, 1250)\n'
    '  [Encoder] 32->64->128->256\n'
    '  [Bottleneck] 512 + FiLM(age,sex)\n'
    '  [Decoder] 256->128->64->32\n'
    'Output: ABP (1, 1250)',
    14, False, BLACK)

add_textbox(slide, 7, 3.8, 5.8, 0.4, 'Why Waveform Reconstruction?', 16, True, PRIMARY)
add_bullet_text(slide, 7, 4.3, 5.8, 2.8, [
    '1. Rich hemodynamic information',
    '   - Not just SBP/DBP numbers',
    '   - Pulse pressure, waveform morphology',
    '   - Dicrotic notch, reflected waves',
    '',
    '2. Clinical interpretation',
    '   - Physicians can visually inspect',
    '   - Detect artifacts and anomalies',
    '',
    '3. Downstream analysis possible',
    '   - Cardiac output estimation',
    '   - Vascular compliance assessment',
], 13, BLACK)


# SLIDE: Final Comprehensive Comparison
slide = add_slide()
add_title_bar(slide, 'Step 4. 전체 모델 비교 (Final)')

add_table_to_slide(slide, 0.3, 1.3, 12.7, 3.5,
    ['Model', 'CV', 'Data', 'SBP MAE', 'DBP MAE', 'MBP MAE', 'SBP R2', 'DBP BHS', 'MBP BHS'],
    [
        ['XGBoost', 'LOSO(46)', '81,812', '14.98', '8.66', '10.41', '-0.19', 'D', 'D'],
        ['LightGBM', 'LOSO(46)', '81,812', '16.09', '8.42', '10.15', '-0.34', 'D', 'D'],
        ['Light CNN (CPU)', '5-Fold', '22,638', '12.95', '7.47', '8.49', '0.16', 'C', 'D'],
        ['Improved CNN (GPU)', '5-Fold', '81,812', '12.64', '7.20', '8.08', '0.17', 'C', 'C'],
        ['ResNet1D (GPU)', '5-Fold', '81,812', '12.25', '7.83', '8.22', '0.20', 'D', 'C'],
        ['U-Net 1D (GPU)', '5-Fold', '81,812', '-', '-', '-', '-', 'Corr=0.68', '-'],
    ])

# Key takeaways
add_rect(slide, 0.3, 5.2, 4, 2, PRGBColor(230, 255, 235))
add_textbox(slide, 0.5, 5.3, 3.6, 0.3, 'Performance Gains', 16, True, GREEN)
add_bullet_text(slide, 0.5, 5.7, 3.6, 1.4, [
    'ML -> CNN: SBP -16%, DBP -17%',
    'Subsample -> Full: DBP 7.47->7.20',
    'CNN -> ResNet: SBP 12.64->12.25',
    'Best DBP: CNN 7.20 (BHS C)',
    'Best SBP: ResNet 12.25 (R2=0.20)',
], 13, BLACK)

add_rect(slide, 4.5, 5.2, 4.2, 2, PRGBColor(255, 245, 230))
add_textbox(slide, 4.7, 5.3, 3.8, 0.3, 'Model Selection Guide', 16, True, PRGBColor(180, 100, 0))
add_bullet_text(slide, 4.7, 5.7, 3.8, 1.4, [
    'Overall best: Improved CNN',
    '  (DBP+MBP Grade C)',
    'Best SBP: ResNet1D',
    '  (MAE 12.25, R2 0.20)',
    'Waveform needed: U-Net',
    '  (Corr 0.68, needs optimization)',
], 13, BLACK)

add_rect(slide, 8.9, 5.2, 4.1, 2, PRGBColor(255, 235, 235))
add_textbox(slide, 9.1, 5.3, 3.7, 0.3, 'Remaining Gaps', 16, True, ACCENT)
add_bullet_text(slide, 9.1, 5.7, 3.7, 1.4, [
    'AAMI not met (MAE>5)',
    'SBP still challenging (>12)',
    'Personal calibration needed',
    'More data would help',
    'U-Net training too slow',
], 13, BLACK)


# SLIDE: Updated Conclusion
slide = add_slide()
add_title_bar(slide, '6. 결론 (Updated with Step 4)')

add_bullet_text(slide, 0.5, 1.4, 12, 5.8, [
    '1.  46명 소아 환자 / 83,439 PPG-ABP 세그먼트 / 6개 모델 비교 완료',
    '',
    '2.  ML Baseline (LOSO): XGBoost SBP MAE=14.98, DBP MAE=8.66 (R2 < 0)',
    '',
    '3.  Improved 1D-CNN (GPU, Full Data):',
    '     -> SBP MAE=12.64, DBP MAE=7.20 (BHS Grade C), MBP MAE=8.08 (BHS Grade C)',
    '     -> Overall best balanced model, recommended as primary choice',
    '',
    '4.  ResNet1D (GPU, Full Data):',
    '     -> SBP MAE=12.25 (best), R2=0.20 (best), skip connections help SBP',
    '',
    '5.  U-Net 1D (Waveform Reconstruction):',
    '     -> PPG -> ABP waveform with Correlation 0.68',
    '     -> FiLM conditioning for age/sex, gradient loss for shape preservation',
    '',
    '6.  Feature importance: sex (21.7%) > age (11.8%) > PPG features',
    '     -> Pediatric demographics are essential for BP prediction',
    '',
    '7.  Next steps: Personal calibration, more data, ensemble, U-Net optimization',
], 16, BLACK)

# Save
pptx_path = os.path.join(BASE, 'PPG2ABP_Presentation.pptx')
try:
    prs.save(pptx_path)
    print("PowerPoint updated.")
except PermissionError:
    alt_path = os.path.join(BASE, 'PPG2ABP_Presentation_v2.pptx')
    prs.save(alt_path)
    print(f"Original PPT locked. Saved as: {alt_path}")

print("\nDone! Both files updated successfully.")
