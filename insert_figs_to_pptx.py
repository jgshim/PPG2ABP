"""Insert pre-generated actual vs predicted figures into PPTX."""
import os
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

BASE = 'C:/Users/jaege/Desktop/Study/PPG2ABP'
FIG_DIR = os.path.join(BASE, 'figures')

# Try to open whichever PPTX is available
for src in ['PPG2ABP_Presentation.pptx', 'PPG2ABP_Presentation_v2.pptx']:
    src_path = os.path.join(BASE, src)
    if os.path.exists(src_path):
        try:
            prs = Presentation(src_path)
            print(f"Opened: {src}")
            break
        except:
            continue

PRIMARY = PRGBColor(0, 51, 102)
WHITE = PRGBColor(255, 255, 255)
GRAY_TEXT = PRGBColor(100, 100, 100)


def add_chart_slide(title, img_path, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, PInches(0), PInches(0), PInches(13.333), PInches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = PRIMARY; shape.line.fill.background()
    txBox = slide.shapes.add_textbox(PInches(0.4), PInches(0.1), PInches(12), PInches(0.7))
    p = txBox.text_frame.paragraphs[0]
    p.text = title; p.font.size = PPt(28); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = 'Malgun Gothic'
    if subtitle:
        slide.shapes.add_picture(img_path, PInches(0.3), PInches(1.05), PInches(12.7), PInches(5.7))
        txBox2 = slide.shapes.add_textbox(PInches(0.3), PInches(6.85), PInches(12.7), PInches(0.4))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = subtitle; p2.font.size = PPt(12); p2.font.color.rgb = GRAY_TEXT
        p2.font.name = 'Malgun Gothic'; p2.alignment = PP_ALIGN.CENTER
    else:
        slide.shapes.add_picture(img_path, PInches(0.3), PInches(1.05), PInches(12.7), PInches(6.2))


def add_two_chart_slide(title, img1, img2, cap1='', cap2=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, PInches(0), PInches(0), PInches(13.333), PInches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = PRIMARY; shape.line.fill.background()
    txBox = slide.shapes.add_textbox(PInches(0.4), PInches(0.1), PInches(12), PInches(0.7))
    p = txBox.text_frame.paragraphs[0]
    p.text = title; p.font.size = PPt(28); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = 'Malgun Gothic'
    slide.shapes.add_picture(img1, PInches(0.2), PInches(1.1), PInches(6.4), PInches(4.8))
    slide.shapes.add_picture(img2, PInches(6.7), PInches(1.1), PInches(6.4), PInches(4.8))
    if cap1:
        txBox1 = slide.shapes.add_textbox(PInches(0.2), PInches(6.0), PInches(6.4), PInches(0.3))
        p1 = txBox1.text_frame.paragraphs[0]; p1.text = cap1
        p1.font.size = PPt(11); p1.font.name = 'Malgun Gothic'
        p1.alignment = PP_ALIGN.CENTER; p1.font.color.rgb = GRAY_TEXT
    if cap2:
        txBox2 = slide.shapes.add_textbox(PInches(6.7), PInches(6.0), PInches(6.4), PInches(0.3))
        p2 = txBox2.text_frame.paragraphs[0]; p2.text = cap2
        p2.font.size = PPt(11); p2.font.name = 'Malgun Gothic'
        p2.alignment = PP_ALIGN.CENTER; p2.font.color.rgb = GRAY_TEXT


# Add slides
add_chart_slide(
    'Actual vs Predicted BP - Time Series (4 Patients)',
    os.path.join(FIG_DIR, 'actual_vs_pred_timeseries.png'),
    'Black=Actual, Red=XGBoost, Blue=CNN, Green=ResNet / Solid=SBP, Dotted=DBP'
)

add_chart_slide(
    'ABP Waveform with Predicted SBP/DBP Overlay',
    os.path.join(FIG_DIR, 'abp_waveform_overlay.png'),
    '5 consecutive 10s segments showing actual ABP waveform with model predictions'
)

add_chart_slide(
    'Single Segment ABP Waveform Detail (4 Patients)',
    os.path.join(FIG_DIR, 'single_segment_detail.png'),
    'Individual 10s ABP waveform with CNN (dashed) and ResNet (dotted) prediction lines'
)

add_two_chart_slide(
    'Actual vs Predicted: Scatter & Error Distribution',
    os.path.join(FIG_DIR, 'scatter_actual_vs_pred.png'),
    os.path.join(FIG_DIR, 'error_distribution.png'),
    'Scatter: closer to diagonal = better prediction',
    'Histogram: narrower and centered at 0 = better'
)

add_chart_slide(
    'Per-Patient MAE Distribution (Box Plot)',
    os.path.join(FIG_DIR, 'per_patient_boxplot.png'),
    'Individual points = per-patient MAE. DL models show lower median and less spread.'
)

# Save
for fname in ['PPG2ABP_Presentation.pptx', 'PPG2ABP_Presentation_v2.pptx', 'PPG2ABP_Presentation_v3.pptx']:
    try:
        save_path = os.path.join(BASE, fname)
        prs.save(save_path)
        print(f"Saved: {save_path}")
        break
    except PermissionError:
        print(f"  {fname} is locked, trying next...")
        continue

print("Done!")
