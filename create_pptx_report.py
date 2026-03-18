"""
Generate a detailed PowerPoint presentation for the PPG2ABP project.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color scheme ──
PRIMARY = RGBColor(0, 51, 102)      # Dark blue
SECONDARY = RGBColor(0, 102, 153)   # Teal blue
ACCENT = RGBColor(220, 80, 60)      # Red
LIGHT_BG = RGBColor(240, 245, 250)  # Light blue-gray
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(40, 40, 40)
GRAY = RGBColor(120, 120, 120)
LIGHT_GRAY = RGBColor(200, 200, 200)
GREEN = RGBColor(40, 160, 80)
TABLE_HEADER = RGBColor(0, 70, 130)
TABLE_ALT = RGBColor(230, 240, 250)


# ── Helper functions ──
def add_slide(layout_idx=6):
    """Add a blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_title_bar(slide, title_text):
    """Add a blue title bar at the top."""
    add_rect(slide, 0, 0, 13.333, 1.1, PRIMARY)
    add_textbox(slide, 0.5, 0.15, 12, 0.8, title_text, 32, True, WHITE, PP_ALIGN.LEFT)

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=BLACK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def add_table_to_slide(slide, left, top, width, height, headers, rows):
    """Add a formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = '맑은 고딕'
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = BLACK
                paragraph.font.name = '맑은 고딕'
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_rect(slide, 0, 0, 13.333, 7.5, PRIMARY)
add_rect(slide, 0, 5.5, 13.333, 0.05, WHITE)

add_textbox(slide, 1, 1.5, 11, 1.5,
            'PPG 파형을 이용한\n침습적 동맥혈압 예측 모델 개발',
            36, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.3, 11, 1,
            'Prediction of Invasive Arterial Blood Pressure\nfrom Photoplethysmography in Pediatric Patients',
            20, False, RGBColor(180, 200, 220), PP_ALIGN.CENTER)
add_textbox(slide, 1, 5.8, 11, 0.8,
            '연구 보고서  |  2026년 3월',
            18, False, RGBColor(180, 200, 220), PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: Outline
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '목차 (Outline)')
items = [
    '1.  서론 — 연구 배경 및 목적',
    '2.  데이터 — 46명 소아 환자, VitalRecorder .vital 파일',
    '3.  방법론',
    '     3-1.  데이터 전처리 (Step 1)',
    '     3-2.  특징 추출 + 머신러닝 (Step 2)',
    '     3-3.  1D-CNN 딥러닝 모델 (Step 3)',
    '4.  결과 — ML vs CNN 비교, 특징 중요도',
    '5.  고찰 — 한계점 및 향후 연구 방향',
    '6.  결론',
]
add_bullet_text(slide, 1.5, 1.5, 10, 5.5, items, 20, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: Background
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '1. 서론 — 연구 배경')

# Left side
add_textbox(slide, 0.5, 1.4, 6, 0.5, '침습적 동맥혈압 (IBP)', 22, True, PRIMARY)
add_bullet_text(slide, 0.5, 2.0, 6, 2.5, [
    '• 가장 정확한 연속 혈압 측정 방법',
    '• 동맥 카테터 삽입 필요 (침습적)',
    '• 합병증 위험: 감염, 혈전, 출혈',
    '• 소아 환자에서 특히 어려움',
], 16, BLACK)

# Right side
add_textbox(slide, 6.8, 1.4, 6, 0.5, 'PPG (광용적맥파)', 22, True, SECONDARY)
add_bullet_text(slide, 6.8, 2.0, 6, 2.5, [
    '• 맥박산소측정기로 비침습적 측정',
    '• 말초 혈관 용적 변화 반영',
    '• 혈관 탄성, 심박출량 정보 포함',
    '• 간편하고 안전한 모니터링',
], 16, BLACK)

add_rect(slide, 1, 4.5, 11.3, 0.05, LIGHT_GRAY)
add_textbox(slide, 0.5, 4.8, 12, 0.5, '핵심 질문', 22, True, ACCENT)
add_textbox(slide, 0.5, 5.4, 12, 1.5,
            'PPG 파형으로부터 침습적 동맥혈압(SBP, DBP, MBP)을 예측할 수 있는가?\n'
            '→ 기존 연구 대부분 성인 대상 — 소아 환자 대상 연구는 매우 부족',
            18, False, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: Objectives
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '1. 서론 — 연구 목적')

objectives = [
    ['1', '소아 환자 PPG → ABP 예측 모델 개발', 'PPG 파형 + 나이/성별을 활용한\n침습적 동맥혈압(SBP, DBP, MBP) 예측'],
    ['2', 'ML vs DL 접근법 비교', 'XGBoost/LightGBM (특징 추출)\nvs 1D-CNN (원시 파형)'],
    ['3', '엄밀한 교차 검증', 'Leave-One-Subject-Out (LOSO)\n5-Fold Group K-Fold'],
    ['4', '임상 표준 기준 평가', 'BHS Grade, AAMI 기준\n(MAE ≤ 5mmHg, SD ≤ 8mmHg)'],
]
for i, (num, title, desc) in enumerate(objectives):
    y = 1.5 + i * 1.4
    add_rect(slide, 0.8, y, 0.7, 0.7, PRIMARY)
    add_textbox(slide, 0.8, y + 0.05, 0.7, 0.6, num, 28, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, 1.8, y, 4, 0.5, title, 18, True, BLACK)
    add_textbox(slide, 6, y, 6.5, 0.7, desc, 14, False, GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: Data Overview
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '2. 데이터 — 개요')

add_table_to_slide(slide, 0.5, 1.4, 5.5, 3.5,
    ['항목', '값'],
    [
        ['환자 수', '46명 (소아)'],
        ['성별', 'M 29명 (63%) / F 17명 (37%)'],
        ['연령 범위', '0.25 ~ 15세'],
        ['데이터 형식', '.vital (VitalRecorder)'],
        ['모니터 장비', 'GE Bx50'],
        ['주요 채널', 'PLETH (PPG), IBP1 (ABP)'],
        ['보조 채널', 'ECG, HR, ART1_SBP/DBP/MBP'],
    ])

add_textbox(slide, 6.5, 1.4, 6, 0.5, '인구통계 정보', 20, True, PRIMARY)
add_bullet_text(slide, 6.5, 2.0, 6, 4, [
    '• PPG_ABP_cases.xlsx에 기록',
    '• 변수: Case ID, Sex (M/F), Age (years)',
    '• 나이 분포: 영아(0.25세) ~ 청소년(15세)',
    '• case13: 3개월 (0.25세) — 최소 연령',
    '',
    '• PPG (PLETH): 광용적맥파 파형',
    '• IBP1: 침습적 동맥혈압 파형 (mmHg)',
    '• ART1_SBP/DBP/MBP: 모니터 자동계산 수치',
    '• 46개 파일 모두 동일한 채널명 사용',
], 14, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: Pipeline Overview
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '3. 방법론 — 전체 파이프라인')

# Step 1 box
add_rect(slide, 0.5, 1.5, 3.8, 2.8, RGBColor(230, 245, 255))
add_textbox(slide, 0.7, 1.6, 3.4, 0.5, 'Step 1: 데이터 전처리', 18, True, PRIMARY)
add_bullet_text(slide, 0.7, 2.2, 3.4, 2, [
    '• .vital 파일 읽기',
    '• 125Hz 리샘플링',
    '• 10초 세그먼트 분할',
    '• 대역통과 필터링',
    '• 품질 검사 & 필터링',
], 13, BLACK)

# Step 2 box
add_rect(slide, 4.7, 1.5, 3.8, 2.8, RGBColor(255, 245, 230))
add_textbox(slide, 4.9, 1.6, 3.4, 0.5, 'Step 2: 특징추출 + ML', 18, True, RGBColor(180, 100, 0))
add_bullet_text(slide, 4.9, 2.2, 3.4, 2, [
    '• 37개 PPG 특징 추출',
    '• 시간/주파수/미분 도메인',
    '• XGBoost, LightGBM',
    '• LOSO 교차검증',
    '• SBP/DBP/MBP 예측',
], 13, BLACK)

# Step 3 box
add_rect(slide, 8.9, 1.5, 3.8, 2.8, RGBColor(230, 255, 235))
add_textbox(slide, 9.1, 1.6, 3.4, 0.5, 'Step 3: 1D-CNN DL', 18, True, RGBColor(0, 120, 60))
add_bullet_text(slide, 9.1, 2.2, 3.4, 2, [
    '• 원시 PPG 파형 입력',
    '• 1D Conv + BN + ReLU',
    '• Late Fusion (age, sex)',
    '• Multi-task (SBP/DBP/MBP)',
    '• 5-Fold Group CV',
], 13, BLACK)

# Bottom summary
add_rect(slide, 0.5, 4.8, 12.3, 2.2, LIGHT_BG)
add_textbox(slide, 0.7, 4.9, 5.5, 0.4, '핵심 설계 원칙', 18, True, PRIMARY)
add_bullet_text(slide, 0.7, 5.4, 11.5, 1.5, [
    '• 환자 단위 교차 검증: 동일 환자의 데이터가 train/test에 겹치지 않도록 엄밀하게 분리',
    '• 생리학적 필터링: SBP 30-250, DBP 10-200, SBP > DBP 조건 적용',
    '• 평가 기준: MAE, RMSE, R², BHS Grade, AAMI 기준 (MAE ≤ 5mmHg, SD ≤ 8mmHg)',
], 14, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: Preprocessing Detail
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '3-1. 데이터 전처리 (Step 1)')

add_table_to_slide(slide, 0.5, 1.4, 12, 4,
    ['단계', '방법', '세부 사항'],
    [
        ['리샘플링', '125Hz 통일', 'vitaldb.read_vital() → to_numpy(interval=1/125)'],
        ['세그먼트 분할', '10초 고정 윈도우', '125Hz × 10s = 1,250 samples/segment'],
        ['품질 필터링', '유효율 ≥ 95%', 'NaN 비율 5% 초과 구간 제외'],
        ['', 'Flat signal 제외', 'PPG 표준편차 < 0.1 인 구간 제외'],
        ['', '생리학적 범위', 'IBP1: 0~300 mmHg 범위 외 제외'],
        ['결측값 처리', '선형 보간', '세그먼트 내 소수 NaN → interpolation'],
        ['PPG 필터', '0.5~8.0 Hz 대역통과', '4차 Butterworth, zero-phase (filtfilt)'],
        ['ABP 필터', '0.5~40.0 Hz 대역통과', '혈압 파형 고조파 보존'],
        ['BP 레이블', '중앙값 (median)', 'ART1_SBP/DBP/MBP의 10초 구간 중앙값'],
    ])

add_rect(slide, 0.5, 5.8, 12, 1.2, LIGHT_BG)
add_textbox(slide, 0.7, 5.9, 11.5, 1,
    '결과: 46명 → 총 83,439 세그먼트 추출 → 필터링 후 81,812 유효 세그먼트\n'
    'SBP 범위: 30~232 mmHg (평균 105±18)  |  DBP 범위: 10~194 mmHg (평균 55±11)',
    16, True, PRIMARY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: Feature Extraction
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '3-2. 특징 추출 + 머신러닝 (Step 2)')

add_textbox(slide, 0.5, 1.3, 6, 0.4, '37개 PPG 특징', 20, True, PRIMARY)
add_table_to_slide(slide, 0.5, 1.8, 6, 4.2,
    ['카테고리', '개수', '특징'],
    [
        ['기본 통계', '6', 'mean, std, skewness, kurtosis, ptp, rms'],
        ['피크 관련', '8', 'n_peaks, HR, ppi_mean/std/cv, peak_amp_mean/std'],
        ['골/진폭', '4', 'valley_amp, pulse_amp 통계'],
        ['시간 특징', '3', 'sys_time, dia_time, sys_dia_ratio'],
        ['면적', '1', 'auc_total'],
        ['미분 (VPG/APG)', '6', 'vpg/apg의 max, min, std'],
        ['주파수 도메인', '7', 'VLF/LF/HF/cardiac power, ratio, entropy'],
        ['인구통계', '2', 'age, sex'],
    ])

add_textbox(slide, 7, 1.3, 5.5, 0.4, '모델 설정', 20, True, PRIMARY)
add_table_to_slide(slide, 7, 1.8, 5.8, 2.5,
    ['파라미터', 'XGBoost', 'LightGBM'],
    [
        ['n_estimators', '300', '300'],
        ['max_depth', '6', '6'],
        ['learning_rate', '0.05', '0.05'],
        ['subsample', '0.8', '0.8'],
        ['reg_alpha / lambda', '0.1 / 1.0', '0.1 / 1.0'],
    ])

add_textbox(slide, 7, 4.6, 5.5, 0.4, '교차 검증', 20, True, PRIMARY)
add_bullet_text(slide, 7, 5.1, 5.5, 2, [
    '• LOSO (Leave-One-Subject-Out): 46 folds',
    '• 각 fold: 1명 테스트, 45명 학습',
    '• StandardScaler: fold별 학습 데이터로 fit',
    '• 가장 엄밀한 일반화 성능 평가',
], 14, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 9: CNN Architecture
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '3-3. 1D-CNN 딥러닝 모델 (Step 3)')

add_textbox(slide, 0.5, 1.3, 6, 0.4, '모델 구조 (PPG2BP_CNN)', 20, True, PRIMARY)
add_table_to_slide(slide, 0.5, 1.8, 6.5, 3.8,
    ['Layer', 'Output', '설명'],
    [
        ['Input', '(1, 1250)', '10초 PPG @ 125Hz'],
        ['Conv1d(1→16, k=15) + BN + ReLU + Pool(4)', '(16, 312)', '넓은 커널로 전체 형태 포착'],
        ['Conv1d(16→32, k=7) + BN + ReLU + Pool(4)', '(32, 78)', '중간 수준 패턴'],
        ['Conv1d(32→64, k=5) + BN + ReLU + Pool(4)', '(64, 19)', '세부 특징'],
        ['Conv1d(64→128, k=3) + BN + ReLU', '(128, 19)', '고수준 표현'],
        ['AdaptiveAvgPool1d(1)', '(128,)', 'Global Average Pooling'],
        ['Concat [age, sex]', '(130,)', 'Late Fusion'],
        ['FC(130→64) + ReLU + Dropout(0.3)', '(64,)', 'Classification head'],
        ['FC(64→3)', '(3,)', 'SBP, DBP, MBP'],
    ])

add_textbox(slide, 7.5, 1.3, 5.5, 0.4, '학습 설정', 20, True, PRIMARY)
add_table_to_slide(slide, 7.5, 1.8, 5.3, 3.2,
    ['항목', '값'],
    [
        ['손실 함수', 'Smooth L1 (Huber)'],
        ['옵티마이저', 'Adam (wd=1e-4)'],
        ['학습률', '2e-3'],
        ['스케줄러', 'ReduceLROnPlateau'],
        ['배치 크기', '256'],
        ['최대 에폭', '20'],
        ['조기 종료', 'Patience=5'],
        ['Grad Clip', 'max_norm=1.0'],
    ])

add_rect(slide, 7.5, 5.5, 5.3, 1.5, LIGHT_BG)
add_textbox(slide, 7.7, 5.6, 5, 0.4, 'Late Fusion 전략', 16, True, PRIMARY)
add_textbox(slide, 7.7, 6.1, 5, 0.8,
    'CNN 특징 추출 후 age/sex를 concat\n→ 파형 패턴과 인구통계 정보를 결합',
    14, False, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 10: ML Results
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '4. 결과 — 머신러닝 모델 (LOSO CV)')

add_table_to_slide(slide, 0.5, 1.4, 12, 3.5,
    ['Target', 'Model', 'MAE (mmHg)', 'RMSE (mmHg)', 'R²', 'Bias±SD (mmHg)', 'BHS Grade', 'AAMI'],
    [
        ['SBP', 'XGBoost', '14.98', '19.47', '-0.1925', '+0.88 ± 19.45', 'D', 'FAIL'],
        ['SBP', 'LightGBM', '16.09', '20.62', '-0.3375', '-0.23 ± 20.62', 'D', 'FAIL'],
        ['DBP', 'XGBoost', '8.66', '11.62', '-0.0780', '+1.76 ± 11.49', 'D', 'FAIL'],
        ['DBP', 'LightGBM', '8.42', '11.33', '-0.0249', '+1.56 ± 11.22', 'D', 'FAIL'],
        ['MBP', 'XGBoost', '10.41', '13.96', '-0.1436', '+1.72 ± 13.85', 'D', 'FAIL'],
        ['MBP', 'LightGBM', '10.15', '13.65', '-0.0935', '+1.43 ± 13.58', 'D', 'FAIL'],
    ])

add_rect(slide, 0.5, 5.3, 12, 1.8, LIGHT_BG)
add_textbox(slide, 0.7, 5.4, 11.5, 0.4, '분석', 18, True, PRIMARY)
add_bullet_text(slide, 0.7, 5.9, 11.5, 1.2, [
    '• R² 음수 → 평균 예측보다 낮은 성능, 환자 간 혈압 변이가 커서 PPG 특징만으로 절대값 예측 어려움',
    '• Bias 작음(< 2 mmHg) → 전체적 편향은 없으나 개별 예측의 분산이 큼',
    '• DBP가 SBP보다 예측 용이 (MAE 8.42 vs 14.98) → 이완기압이 PPG와 더 밀접한 관계',
], 14, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 11: CNN Results
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '4. 결과 — 1D-CNN 모델 (5-Fold Group CV)')

add_table_to_slide(slide, 0.5, 1.4, 12, 2.5,
    ['Target', 'MAE (mmHg)', 'RMSE (mmHg)', 'R²', 'Bias±SD', '≤5mmHg', '≤10mmHg', '≤15mmHg', 'BHS'],
    [
        ['SBP', '12.95', '16.76', '0.1569', '+0.5 ± 16.8', '25.8%', '47.8%', '65.4%', 'D'],
        ['DBP', '7.47', '10.21', '0.2046', '-0.1 ± 10.2', '42.8%', '74.8%', '90.6%', 'C'],
        ['MBP', '8.49', '11.67', '0.2193', '-0.1 ± 11.7', '39.1%', '68.8%', '84.4%', 'D'],
    ])

add_textbox(slide, 0.5, 4.2, 6, 0.4, 'Fold별 상세 결과', 18, True, PRIMARY)
add_table_to_slide(slide, 0.5, 4.7, 6, 2.2,
    ['Fold', 'Test Subj', 'Samples', 'SBP MAE', 'DBP MAE', 'MBP MAE'],
    [
        ['1', '9명', '4,500', '15.06', '7.56', '9.24'],
        ['2', '9명', '4,500', '13.24', '6.94', '8.41'],
        ['3', '9명', '4,500', '10.37', '7.69', '7.52'],
        ['4', '9명', '4,500', '13.00', '7.69', '8.69'],
        ['5', '10명', '4,638', '13.09', '7.49', '8.60'],
    ])

add_rect(slide, 7, 4.2, 5.8, 2.8, RGBColor(230, 255, 235))
add_textbox(slide, 7.2, 4.3, 5.4, 0.4, 'Key Findings', 18, True, GREEN)
add_bullet_text(slide, 7.2, 4.8, 5.4, 2.2, [
    '• DBP: BHS Grade C 달성!',
    '  ≤5mmHg 42.8% (기준 40%)',
    '  ≤10mmHg 74.8% (기준 65%)',
    '  ≤15mmHg 90.6% (기준 85%)',
    '',
    '• R² 양수 전환 → CNN이 환자 간',
    '  변이를 일부 설명 가능',
    '• Bias ≈ 0 → 편향 없는 예측',
], 14, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 12: Comparison
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '4. 결과 — 모델 비교 (ML vs CNN)')

add_table_to_slide(slide, 0.8, 1.5, 11.5, 3.5,
    ['Target', 'Model', 'CV 방법', 'MAE (mmHg)', 'RMSE (mmHg)', 'R²', '개선율'],
    [
        ['SBP', 'XGBoost', 'LOSO (46-fold)', '14.98', '19.47', '-0.1925', '-'],
        ['SBP', '1D-CNN', 'Group 5-fold', '12.95', '16.76', '+0.1569', '↓13.5%'],
        ['DBP', 'XGBoost', 'LOSO (46-fold)', '8.66', '11.62', '-0.0780', '-'],
        ['DBP', '1D-CNN', 'Group 5-fold', '7.47', '10.21', '+0.2046', '↓13.7%'],
        ['MBP', 'XGBoost', 'LOSO (46-fold)', '10.41', '13.96', '-0.1436', '-'],
        ['MBP', '1D-CNN', 'Group 5-fold', '8.49', '11.67', '+0.2193', '↓18.4%'],
    ])

add_rect(slide, 0.8, 5.5, 5.5, 1.5, LIGHT_BG)
add_textbox(slide, 1, 5.6, 5.3, 0.4, 'CNN 우위 요인', 16, True, PRIMARY)
add_bullet_text(slide, 1, 6.0, 5, 1, [
    '• 원시 파형에서 자동 특징 학습',
    '• 수작업 특징으로 포착 못한 패턴 학습',
    '• Multi-task 학습의 정규화 효과',
], 13, BLACK)

add_rect(slide, 6.8, 5.5, 5.7, 1.5, RGBColor(255, 240, 240))
add_textbox(slide, 7, 5.6, 5.3, 0.4, '주의사항', 16, True, ACCENT)
add_bullet_text(slide, 7, 6.0, 5.3, 1, [
    '• CV 방법이 다름 (LOSO vs 5-fold)',
    '• CNN은 서브샘플링 사용 (22,638 샘플)',
    '• 공정 비교를 위해 동일 CV 필요',
], 13, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 13: Feature Importance
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '4. 결과 — 특징 중요도 분석 (XGBoost, SBP)')

add_table_to_slide(slide, 0.5, 1.4, 6, 5,
    ['순위', '특징명', '중요도', '카테고리'],
    [
        ['1', 'sex (성별)', '0.2166', '인구통계'],
        ['2', 'age (나이)', '0.1177', '인구통계'],
        ['3', 'dominant_freq', '0.0546', '주파수'],
        ['4', 'peak_amp_mean', '0.0544', '피크'],
        ['5', 'apg_min', '0.0445', '미분(APG)'],
        ['6', 'skewness', '0.0402', '통계'],
        ['7', 'spectral_entropy', '0.0395', '주파수'],
        ['8', 'ppi_cv', '0.0360', '피크'],
        ['9', 'sys_time_mean', '0.0335', '시간'],
        ['10', 'apg_std', '0.0323', '미분(APG)'],
    ])

add_rect(slide, 7, 1.4, 5.8, 2.5, LIGHT_BG)
add_textbox(slide, 7.2, 1.5, 5.4, 0.4, '핵심 인사이트', 18, True, PRIMARY)
add_bullet_text(slide, 7.2, 2.0, 5.4, 1.8, [
    '• 성별(21.7%)과 나이(11.8%)가 압도적 1, 2위',
    '• 소아 환자에서 연령별 정상 혈압이',
    '  크게 다르기 때문 (예: 1세 vs 15세)',
    '• 인구통계 정보가 PPG 특징보다',
    '  혈압 예측에 더 중요한 역할',
], 14, BLACK)

add_rect(slide, 7, 4.2, 5.8, 2.5, RGBColor(255, 250, 230))
add_textbox(slide, 7.2, 4.3, 5.4, 0.4, 'PPG 특징 중에서는', 18, True, RGBColor(180, 100, 0))
add_bullet_text(slide, 7.2, 4.8, 5.4, 1.8, [
    '• dominant_freq: 심박수 관련',
    '• peak_amp_mean: 맥파 진폭',
    '• apg_min: 2차 미분 최솟값',
    '  (동맥 경직도와 관련)',
    '• spectral_entropy: 파형 복잡도',
    '• sys_time_mean: 수축기 시간',
], 14, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 14: Discussion
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '5. 고찰 — 한계점 및 향후 연구')

add_textbox(slide, 0.5, 1.3, 6, 0.4, '현재 한계점', 20, True, ACCENT)
add_bullet_text(slide, 0.5, 1.9, 6, 3, [
    '1. 데이터 규모: 46명 — 더 많은 환자 필요',
    '2. 개인 보정 부재: 절대 혈압 예측의 한계',
    '3. 수술 중 데이터: 마취/약물 영향',
    '4. AAMI 기준 미달: MAE > 5mmHg',
    '5. 파형 재구성 미수행',
], 15, BLACK)

add_textbox(slide, 6.8, 1.3, 6, 0.4, '향후 연구 방향', 20, True, GREEN)
add_bullet_text(slide, 6.8, 1.9, 6, 3, [
    '1. 개인 보정 (Calibration)',
    '   → 초기 NIBP로 offset 보정 (최대 개선 기대)',
    '2. U-Net 파형 재구성',
    '   → PPG → ABP 파형 전체 생성',
    '3. ResNet1D / Transformer',
    '   → GPU 활용 대규모 학습',
    '4. 다중 신호 (ECG + PPG)',
    '5. 앙상블 (ML + CNN 결합)',
    '6. 외부 검증 (다기관 데이터)',
], 14, BLACK)

add_rect(slide, 0.5, 5.3, 12, 1.8, RGBColor(230, 255, 235))
add_textbox(slide, 0.7, 5.4, 11.5, 0.4, 'GPU 환경 구축 완료', 18, True, GREEN)
add_textbox(slide, 0.7, 5.9, 11.5, 1,
    '• NVIDIA GeForce RTX 4070 Laptop (8GB VRAM)\n'
    '• PyTorch 2.6.0 + CUDA 12.4 설치 완료\n'
    '• 전체 데이터(81,812 샘플)로 대규모 모델 학습 가능 → LOSO 46-fold CNN도 수 분 내 완료 예상',
    14, False, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 15: Conclusion
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_title_bar(slide, '6. 결론')

add_bullet_text(slide, 0.8, 1.5, 11.5, 5.5, [
    '1.  46명 소아 환자에서 83,439개의 PPG-ABP 세그먼트 쌍을 추출하였다.',
    '',
    '2.  XGBoost (LOSO): SBP MAE 14.98 mmHg, DBP MAE 8.66 mmHg',
    '     → PPG 특징만으로는 환자 간 절대 혈압 예측에 한계 (R² < 0)',
    '',
    '3.  1D-CNN (5-Fold Group CV): SBP MAE 12.95 mmHg, DBP MAE 7.47 mmHg',
    '     → ML 대비 13~18% 성능 향상, R² 양수 전환',
    '     → DBP 예측에서 BHS Grade C 달성',
    '',
    '4.  특징 중요도: 성별(21.7%) > 나이(11.8%) > PPG 파형 특징',
    '     → 소아 환자 혈압 예측에서 인구통계 정보의 필수적 역할 확인',
    '',
    '5.  임상 표준(AAMI)에는 미달하나, 소아 PPG→ABP 연구의 기초를 마련',
    '     → 개인 보정, 더 큰 모델, 더 많은 데이터로 성능 향상 기대',
], 17, BLACK)


# ════════════════════════════════════════════════════════════════
# SLIDE 16: Thank you
# ════════════════════════════════════════════════════════════════
slide = add_slide()
add_rect(slide, 0, 0, 13.333, 7.5, PRIMARY)
add_textbox(slide, 1, 2.5, 11, 1, 'Thank You', 48, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1, 4, 11, 1,
            'Questions & Discussion',
            24, False, RGBColor(180, 200, 220), PP_ALIGN.CENTER)


# ── Save ──
output_path = 'C:/Users/jaege/Desktop/Study/PPG2ABP/PPG2ABP_Presentation.pptx'
prs.save(output_path)
print(f"PowerPoint saved: {output_path}")
